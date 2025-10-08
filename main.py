import sys
import os
import webbrowser
from io import BytesIO
import math

from PyQt5.QtWidgets import (
    QApplication, QWidget, QLabel, QLineEdit, QDialog, QPushButton, 
    QVBoxLayout, QHBoxLayout, QGroupBox, QRadioButton,  QFileDialog,
    QDateEdit, QToolButton, QComboBox, QDoubleSpinBox, QGraphicsView,
    QGraphicsScene, QGraphicsPixmapItem, QGraphicsProxyWidget, QFrame, 
    QToolBar, QGridLayout, QGroupBox, QCheckBox, QTableWidget, QTableWidgetItem,
    QSizePolicy, QMessageBox, QListWidgetItem, QScrollArea, 
    QListWidget, QMenu,
)
from PyQt5.QtPrintSupport import QPrintDialog, QPrintPreviewWidget, QPrinter, QPrintPreviewDialog
from PyQt5.QtGui import QPixmap, QImage, QPainter, QPageSize, QPen, QColor, QTextDocument, QCursor
from PyQt5.QtCore import Qt, QUrl, QSize, QSizeF, QDate, QPointF, QRectF

MEDIA_ROOT = os.path.join(os.getcwd(), 'media')  # The main folder
TEMPLATE_DIR = os.path.join(MEDIA_ROOT, 'templates') # For blank shirt images (ComboBox source)
REFERENCE_DIR = os.path.join(MEDIA_ROOT, 'references') # For customer-uploaded photos (Gallery source) 

class ImageGalleryWindow(QDialog):
    def __init__(self, media_dir, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Uploaded Image Gallery")
        self.media_dir = media_dir  # Store the path for load_images()
        self.setGeometry(100, 100, 800, 600)
        
        main_layout = QVBoxLayout(self)
        
        scroll_area = QScrollArea()
        scroll_area.setWidgetResizable(True)
        main_layout.addWidget(scroll_area)
        
        self.gallery_widget = QWidget()
        self.gallery_layout = QGridLayout(self.gallery_widget)
        scroll_area.setWidget(self.gallery_widget)
        
        self.load_images()
    
    def load_images(self):
        image_extensions = ('.png', '.jpg', '.jpeg', '.bmp', '.webp')
        col_count = 4 
        row, col = 0, 0

        if not os.path.isdir(self.media_dir):
            self.gallery_layout.addWidget(QLabel("Error: Reference directory not found."), 0, 0)
            return

        for filename in os.listdir(self.media_dir):
            if filename.lower().endswith(image_extensions):
                image_path = os.path.join(self.media_dir, filename)
                
                label = QLabel()
                pixmap = QPixmap(image_path)
                
                if not pixmap.isNull():
                    label.setPixmap(pixmap.scaled(200, 200, Qt.KeepAspectRatio, Qt.SmoothTransformation))
                    label.setAlignment(Qt.AlignCenter)
                    label.setToolTip(filename) # Show filename on hover
                    
                    self.gallery_layout.addWidget(label, row, col)
                    
                    col += 1
                    if col >= col_count:
                        col = 0
                        row += 1

class OrderForm(QWidget):
    def __init__(self):
        super().__init__()
        self.setWindowTitle("Order Form")
        self.setStyleSheet("background-color: #F5FFFA")

    
        # ✅ सबसे पहले image और बाकी variables init करो
        self.image = None
        self.scene = None
        self.canvas = None
        self.entries = {}
        self._canvas_image_item = None

        # 🔹 Main vertical layout
        self.main_layout = QVBoxLayout(self)
        self.main_layout.setContentsMargins(5, 5, 5, 5)
        self.main_layout.setSpacing(2)

        # Toolbar
        toolbar = self.main_toolbar()
        self.main_layout.addWidget(toolbar, alignment=Qt.AlignLeft | Qt.AlignTop)

        # Header
        header = self.header()
        self.main_layout.addWidget(header, alignment=Qt.AlignLeft | Qt.AlignTop)

        # Product + Image Panel
        product_panel = self._product_and_image_panel()
        self.main_layout.addWidget(product_panel, alignment=Qt.AlignLeft | Qt.AlignTop)

       
        # Pehla row add karte hi dikhega
        item_box=self.create_item_selection_box()
        self.main_layout.addWidget(item_box)

        self.add_remark_field()

      
        self.main_layout.addLayout(self.create_buttons_row())
        # Add row panel first


# Then stretch at the very end
        self.main_layout.addStretch()

    def create_button(self, text, emoji, shortcut=None, size=(80, 80)):
        """ Helper to create styled buttons with emoji as icon (top) """
        btn = QToolButton()
        btn.setText(f"{emoji}\n{text}")   # Icon ऊपर, text नीचे
        btn.setToolButtonStyle(Qt.ToolButtonTextUnderIcon)
        btn.setFixedSize(*size)

        # Shortcut
        if shortcut:
            btn.setShortcut(shortcut)

        # Styling
        btn.setStyleSheet("""
            QToolButton {
                background-color: #EEEEEE;
                color: #fc83a0;
                font-weight: bold;
                border: 1px solid #CCCCCC;
            }
            QToolButton:hover {
                background-color: #fc83a0;
                color: white;
            }
        """)
        return btn

    def main_toolbar(self):
        """
        Top panel with action buttons (NEW, EDIT, DELETE, SEARCH, etc.)
        """
        button_group = QWidget()

        # 🔹 Main horizontal layout
        button_layout = QHBoxLayout()
        button_layout.setContentsMargins(5, 5, 5, 5)
        button_layout.setSpacing(0)

        # 🔹 Sub-layout 1 (NEW/EDIT/DELETE/SEARCH – chipake)
        left_layout = QHBoxLayout()
        left_layout.setSpacing(1)
        left_layout.setContentsMargins(0, 0, 0, 0)

        # 🔹 Sub-layout 2 (TOP/BACK/NEXT/LAST – chipake)
        mid_layout = QHBoxLayout()
        mid_layout.setSpacing(1)
        mid_layout.setContentsMargins(0, 0, 0, 0)

        # 🔹 Sub-layout 3 (EXIT/TUTOR – chipake)
        right_layout = QHBoxLayout()
        right_layout.setSpacing(1)
        right_layout.setContentsMargins(0, 0, 0, 0)

        # Buttons with shortcuts
        self.new_btn = QPushButton("📄\n New-Ctrl+N");    self.new_btn.setShortcut("Ctrl+N")
        self.edit_btn = QPushButton("📝\n Modify-Ctrl+O");  self.edit_btn.setShortcut("Ctrl+E")
        self.delete_btn = QPushButton("🗑️\n Delete"); self.delete_btn.setShortcut("Ctrl+D")
        self.search_btn = QPushButton("🔍\n Search-Ctrl+F"); self.search_btn.setShortcut("Ctrl+F")
        self.search_btn.clicked.connect(self.open_search_window)
        self.print_btn = QPushButton("🖨\n Print-Ctrl+P"); self.print_btn.setShortcut("Ctrl+P")

        self.top_btn = QPushButton("▲\n Top")
        self.back_btn = QPushButton("◀\n Back");  self.back_btn.setShortcut("Alt+Left")
        self.next_btn = QPushButton("▶\n Next");  self.next_btn.setShortcut("Alt+Right")
        self.last_btn = QPushButton("▼\n Last")

        self.exit_btn = QPushButton("↩️\n Exit-Alt+F4");  self.exit_btn.setShortcut("Alt+F4")
        self.tutor_btn = QPushButton("❔\n TUTOR")
        self.upload_btn=QPushButton("Upload")
        self.previous_btn=QPushButton("previous")

       #Connections 

        self.upload_btn.clicked.connect(self._upload_reference_image) 
        self.previous_btn.clicked.connect(self._open_image_gallery) 
        self.print_btn.clicked.connect(self.open_print_dialog)
       # Common style
        button_style = """
            background-color: #EEEEEE; 
            color: #fc83a0; 
            font-weight: bold; 
            padding: 10px;
        """
        all_buttons = [
            self.new_btn, self.edit_btn, self.delete_btn, self.search_btn, self.print_btn,
            self.top_btn, self.back_btn, self.next_btn, self.last_btn,
            self.exit_btn, self.tutor_btn,self.upload_btn,self.previous_btn
        ]
        for btn in all_buttons:
            btn.setStyleSheet(button_style)

        # 🔹 Group 1 (left buttons)
        for btn in [self.new_btn, self.edit_btn, self.delete_btn, self.search_btn, self.print_btn]:
            btn.setFixedSize(150, 70)
            left_layout.addWidget(btn)

        # 🔹 Group 2 (mid buttons)
        for btn in [self.top_btn, self.back_btn, self.next_btn, self.last_btn]:
            btn.setFixedSize(80, 70)
            mid_layout.addWidget(btn)

        # 🔹 Group 3 (right buttons)
        for btn in [self.exit_btn, self.tutor_btn,self.upload_btn,self.previous_btn]:
            # btn.setFixedSize(150, 70)
            # right_layout.addWidget(btn)
            if btn in [self.upload_btn, self.previous_btn]:
                btn.setFixedSize(150, 40)  # height kam
            else:
                btn.setFixedSize(150, 70)  # normal height
            right_layout.addWidget(btn)

            if btn==self.tutor_btn:
                right_layout.addSpacing(20)

        # 🔹 Add layouts to main layout with spacing only between groups
        button_layout.addLayout(left_layout)
        button_layout.addSpacing(20)   # left & mid gap
        button_layout.addLayout(mid_layout)
        button_layout.addSpacing(20)   # mid & right gap
        button_layout.addLayout(right_layout)

        button_group.setLayout(button_layout)
        return button_group

    def header(self, parties=None):
        container = QWidget()
        main_layout = QVBoxLayout(container)

        button_style = """
            background-color: #EEEEEE; 
            color: #fc83a0; 
            font-weight: bold; 
            padding: 10px;
        """

        # ---------- 1st Row ----------
        row1 = QHBoxLayout()
        self.order_number = QLineEdit()
        self.order_number.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
      
        self.order_date = QDateEdit()
        self.order_date.setCalendarPopup(True)
        self.order_date.setDisplayFormat("dd-MM-yyyy")
        self.order_date.setDate(QDate.currentDate())
        self.order_date.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        
       
        self.delivery_date = QDateEdit()
        self.delivery_date.setCalendarPopup(True)
        self.delivery_date.setDisplayFormat("dd-MM-yyyy")
        self.delivery_date.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        
        
        self.barcode_number = QLineEdit()
        self.barcode_number.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        self.save_barcode_btn = QPushButton("Save Barcode")
        self.save_barcode_btn.setStyleSheet(button_style)
        
        self.gst_no = QLineEdit()
        self.gst_no.setFixedSize(200,30)
        self.gst_no.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        
        self.advance_paid = QDoubleSpinBox()
        self.advance_paid.setPrefix("₹ ")
        self.advance_paid.setMaximum(9999999)
        self.advance_paid.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")


        row1.addWidget(QLabel("Order No:"))
        row1.addWidget(self.order_number)
        row1.addWidget(QLabel("Order Date:"))
        row1.addWidget(self.order_date)
        row1.addWidget(QLabel("Delivery Date:"))
        row1.addWidget(self.delivery_date)
        row1.addWidget(QLabel("Barcode:"))
        row1.addWidget(self.barcode_number)
        row1.addWidget(self.save_barcode_btn)
        row1.addWidget(QLabel("GST No:"))
        row1.addWidget(self.gst_no)
        row1.addWidget(QLabel("Advance Paid:"))
        row1.addWidget(self.advance_paid)

        # ---------- 2nd Row ----------
        row2 = QHBoxLayout()
        self.party_name = QLineEdit()
        if parties:
            self.party_name.addItems(parties)
        self.party_name.setFixedSize(250,30)
        self.party_name.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        self.school_name = QLineEdit()
        self.school_name.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        # self.gst_no = QLineEdit()
        # self.gst_no.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        self.address = QLineEdit()
        self.address.setStyleSheet("background-color: white; border: 1px solid gray; padding: 3px;")
        self.address.setFixedHeight(30)

        row2.addWidget(QLabel("Party Name:"))
        row2.addWidget(self.party_name)
        row2.addWidget(QLabel("School Name:"))
        row2.addWidget(self.school_name)
        # row2.addWidget(QLabel("GST No:"))
        # row2.addWidget(self.gst_no)
        row2.addWidget(QLabel("Address:"))
        row2.addWidget(self.address)
        
        # Add all rows
        main_layout.addLayout(row1)
        main_layout.addLayout(row2)

        return container
    
    def _open_image_gallery(self):
        
        gallery = ImageGalleryWindow(media_dir=REFERENCE_DIR, parent=self)
        gallery.exec_() 

    def _upload_reference_image(self):
        
        path, _ = QFileDialog.getOpenFileName(
            self, "Select Reference Image", "", "Images (*.png *.jpg *.jpeg *.bmp *.webp)"
        )
        
        if path:
            temp_pixmap = QPixmap(path)            
            original_filename = os.path.basename(path)
            destination_path = os.path.join(REFERENCE_DIR, original_filename)             
            if temp_pixmap.isNull():
                print("Error: Could not load the image.")
                return

            if temp_pixmap.save(destination_path):
                print(f"Customer reference image saved to: {destination_path}")
                
            else:
                print(f"Error: Could not save reference image to {destination_path}")
                        
    def _product_and_image_panel(self):
        """Middle row: left = product image canvas, right = options card"""
        container = QWidget()
        main_layout = QHBoxLayout(container)
        main_layout.setSpacing(30)  

        # LEFT PANEL
        left_layout = QVBoxLayout()
        top_bar = QHBoxLayout()

        lbl = QLabel("PRODUCT IMAGE :")
        lbl.setStyleSheet("border: none;") 
        lbl.setFixedWidth(170)
        self.cmb = QComboBox()
        self.display_to_filename_map = {}
        self.display_names = ["SELECT"] 
        if os.path.isdir(TEMPLATE_DIR):
            image_extensions = ('.png', '.jpg', '.jpeg', '.bmp', '.webp')
            
            for filename in os.listdir(TEMPLATE_DIR):
                if filename.lower().endswith(image_extensions):
                    name_without_ext = os.path.splitext(filename)[0]
                    display_name = name_without_ext.upper()                    
                    self.display_names.append(display_name)
                    self.display_to_filename_map[display_name] = filename
        
        self.cmb.addItems(self.display_names)

        self.cmb.setFixedWidth(200)

        self.cmb.currentTextChanged.connect(self._change_image_from_select)

        btn_upload = QPushButton("+")
        btn_upload.setFixedWidth(80)
        btn_upload.setStyleSheet("background-color:#2563eb; color:white;")
        btn_upload.clicked.connect(self._upload_image)

        top_bar.addWidget(lbl)
        top_bar.addWidget(self.cmb)
        top_bar.addWidget(btn_upload)

        # Graphics canvas
        self.scene = QGraphicsScene()
        self.canvas = QGraphicsView(self.scene)
        self.canvas.setFixedSize(525, 300)
        self.canvas.setStyleSheet("background:#fafafa; border:1px solid #d1d5db;")

        left_layout.addLayout(top_bar)
        left_layout.addWidget(self.canvas)
        left_frame = QFrame()   # 👈 अब define कर दिया
        left_frame.setLayout(left_layout)
        left_frame.setFixedWidth(550)
        left_frame.setStyleSheet("background:white; border:1px solid #ddd;")
     
        # 🔹 RIGHT PANEL (Options)
       
        right_frame = QFrame()
        right_layout = QVBoxLayout(right_frame)
        build_option = self._build_options_panel()   # 👈 अब parent नको
        right_layout.addWidget(build_option)
        #right_frame.setFixedWidth(700)   # 👈 इथे width adjust करा
        # right_frame.setStyleSheet("background:white; border:1px solid #ddd;")

        main_layout.addWidget(left_frame)
        # main_layout.addWidget(second_frame)
        main_layout.addWidget(right_frame)
        default_file = "NEW REGULAR COLER.jpg"
        default_image_path = os.path.join(TEMPLATE_DIR, default_file)

        if os.path.exists(default_image_path):
            self.image = QPixmap(default_image_path)
            
            default_display_name = os.path.splitext(default_file)[0].upper()
            
            index = self.cmb.findText(default_display_name)
            if index != -1:
                self.cmb.setCurrentIndex(index)
                
        else:
            self.image = None
            self.cmb.setCurrentText("SELECT")
        self._render_image()

        return container
    
    def _change_image_from_select(self, display_name):
        if display_name == "SELECT" or not display_name:
            return
        filename = self.display_to_filename_map.get(display_name)
        if filename:
            image_path = os.path.join(TEMPLATE_DIR, filename)
            
            if os.path.exists(image_path):
                new_image = QPixmap(image_path)
                
                if not new_image.isNull():
                    self.image = new_image
                    self._render_image()

    def _build_options_panel(self):
        parent = QWidget()
        layout = QGridLayout(parent)

       
        INPUT_STYLE = """
            background-color: white;
            border: 1px solid gray;
            padding: 3px;
        """

    # ========== Printing Options ==========
        default_prices = {'front': '5', 'back': '7', 'patch': '5', 'embroidery': '15', 'Dtf': '0', 'Front sablimation': '60', 'Back sablimation': '60'}
        sec_print = QGroupBox("Printing Options", parent)
        grid_print = QGridLayout(sec_print)
        sec_print.setMaximumWidth(350)
        sec_print.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Preferred)

        sec_print.setContentsMargins(30,0,20,0)
        self.print_vars = {}
        keys = ['front', 'back', 'patch', 'embroidery', 'Dtf', 'Front sablimation', 'Back sablimation']

        for idx, key in enumerate(keys):
            default_price = default_prices.get(key, "0")
            cb = QCheckBox(key.upper())
            price_edit = QLineEdit(default_price)
            price_edit.setStyleSheet(INPUT_STYLE)
            price_edit.setEnabled(False)

            def toggle(state, entry=price_edit):
                entry.setEnabled(state == Qt.Checked)
                #if state != Qt.Checked:
                #    entry.setText("0.0")

            cb.stateChanged.connect(toggle)
            self.print_vars[key] = (cb, price_edit)

            grid_print.addWidget(cb, idx, 0)
            grid_print.addWidget(QLabel("Price"), idx, 1)
            grid_print.addWidget(price_edit, idx, 2)

        layout.addWidget(sec_print, 0, 0)

   # ========== Collar Options ==========
        sec1 = QGroupBox("Collar Options", parent)
        grid1 = QGridLayout(sec1)
        sec1.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Preferred)
        sec1.setMinimumWidth(250)
        sec1.setContentsMargins(20,0,50,0)
        sec1.move(sec1.x(), 10)  # y=10 px upar se


        self.collar_var = "self"
        self.collar_price_self = QLineEdit("0")
        self.collar_price_self.setStyleSheet(INPUT_STYLE)
        self.collar_price_rib = QLineEdit("10")
        self.collar_price_rib.setStyleSheet(INPUT_STYLE)
        self.collar_price_patti = QLineEdit("10")
        self.collar_price_patti.setStyleSheet(INPUT_STYLE)
        self.collar_cloth = QComboBox()
        self.collar_cloth.addItems(["Cotton", "Polyester", "Blended", "Other"])

        rb_self = QCheckBox("Self Collar")
        rb_self.setChecked(True)
        rb_self.toggled.connect(lambda checked: setattr(self, "collar_var", "self") if checked else None)
        grid1.addWidget(rb_self, 0, 0)
        grid1.addWidget(QLabel("Price"), 0, 1)
        grid1.addWidget(self.collar_price_self, 0, 2)

        rb_rib = QCheckBox("RIB collar")
        rb_rib.toggled.connect(lambda checked: setattr(self, "collar_var", "rib") if checked else None)
        grid1.addWidget(rb_rib, 1, 0)
        grid1.addWidget(QLabel("Price"), 1, 1)
        grid1.addWidget(self.collar_price_rib, 1, 2)

        rb_patti = QCheckBox("RIB Patti")
        rb_patti.toggled.connect(lambda checked: setattr(self, "collar_var", "patti") if checked else None)
        grid1.addWidget(rb_patti, 2, 0)
        grid1.addWidget(QLabel("Price"), 2, 1)
        grid1.addWidget(self.collar_price_patti, 2, 2)

        """other_patti = QCheckBox("Other Patti")
        other_patti.toggled.connect(lambda checked: setattr(self, "collar_var", "patti") if checked else None)
        grid1.addWidget(rb_patti, 2, 0)
        grid1.addWidget(QLabel("Price"), 2, 1)
        grid1.addWidget(self.collar_price_patti, 2, 2)"""
        layout.addWidget(sec1, 0, 1)
 
        # ========== Button and Style Options ==========
        sec_button = QGroupBox("Button Options", parent)
        grid_button = QGridLayout(sec_button)

        self.style_var = "button"
        rb_button = QCheckBox("BUTTON")
        rb_button.setFixedWidth(200)
        rb_button.setChecked(True)
        rb_button.toggled.connect(lambda: setattr(self, "style_var", "button"))
        grid_button.addWidget(rb_button, 0, 0)

        rb_plain = QRadioButton("PLAIN")
        # rb_plain.setChecked(True)
        rb_plain.toggled.connect(lambda: setattr(self, "style_var", "plain"))
        grid_button.addWidget(rb_plain, 1, 0)

        rb_box = QRadioButton("BOX")
        rb_box.toggled.connect(lambda: setattr(self, "style_var", "box"))
        grid_button.addWidget(rb_box, 2, 0)

        rb_vplus = QRadioButton("V+")
        rb_vplus.toggled.connect(lambda: setattr(self, "style_var", "vplus"))
        grid_button.addWidget(rb_vplus, 3, 0)

        layout.addWidget(sec_button,0,2)
        layout.setColumnStretch(0, 0)
        layout.setColumnStretch(1, 0)       
        layout.setColumnStretch(2, 0)        
        layout.setColumnStretch(3, 1)
       
        # Status layout
        COMBO_BOX_STYLE = """
            QComboBox {
                /* Define a clean, solid border */
                border: 1px solid #777; 
                border-radius: 3px;
            }
        """

        status_widget = QWidget()
        status_layout = QHBoxLayout(status_widget)
        status_layout.setContentsMargins(2,2,2,2)
        status_layout.setSpacing(10) 

        status_label = QLabel("Status:")
        status_label.setMinimumWidth(50)
        font = status_label.font()
        font.setPointSize(11) 
        status_label.setFont(font)
        status_layout.addWidget(status_label)

        # Main status combo
        self.status_combo = QComboBox()
        self.status_combo.addItem("Pending")
        self.status_combo.addItem("Running")
        self.status_combo.addItem("Completed")
        self.status_combo.setFixedWidth(150)
        self.status_combo.setStyleSheet(COMBO_BOX_STYLE) 
        font = self.status_combo.font()
        font.setPointSize(11)
        self.status_combo.setFont(font)
        status_layout.addWidget(self.status_combo)

        # Sub-options combo (initially hidden)
        self.sub_combo = QComboBox()
        self.sub_combo.addItems(["Cutting", "Streching", "Printing"])
        self.sub_combo.setFixedWidth(200)
        self.sub_combo.setStyleSheet(COMBO_BOX_STYLE) 
        self.sub_combo.setVisible(False)
        font = self.sub_combo.font()
        font.setPointSize(11)
        self.sub_combo.setFont(font)
        status_layout.addWidget(self.sub_combo)

        status_layout.addStretch()

        # Connect main combo to show/hide sub-combo
        self.status_combo.currentTextChanged.connect(
            lambda value: self.sub_combo.setVisible(value == "Running")
        )

        # Sub-options list
        layout.addWidget(status_widget, 1, 1, 1, 2)

        # ===== Finally return parent =====
        return parent
                
    def _upload_image(self):
        path, _ = QFileDialog.getOpenFileName(
            self, "Select Image", "", "Images (*.png *.jpg *.jpeg *.bmp *.webp)"
        )
        if path:
            self.image = QPixmap(path)
            self._render_image()

        # 3. Define the new filename for storage
        # We use the original file's name but ensure it's saved in MEDIA_DIR
        original_filename = os.path.basename(path)
        
        # Create a display name (e.g., "my_shirt.png" -> "MY_SHIRT")
        name_without_ext = os.path.splitext(original_filename)[0]
        display_name = name_without_ext.upper()
        
        # 4. Save the file permanently to MEDIA_DIR
        destination_path = os.path.join(TEMPLATE_DIR, original_filename)
        
        # Only save if the file isn't already there (or overwrite it)
        if not os.path.exists(destination_path) or path != destination_path:
            # We use the QPixmap's save method for reliable saving
            # NOTE: self.image is the QPixmap loaded from 'path'
            success = self.image.save(destination_path)
            
            if success:
                print(f"Image saved to: {destination_path}")
            else:
                print(f"Error: Could not save image to {destination_path}")
                return # Exit if saving failed

        # 5. Update the ComboBox and mapping dictionary
        if display_name not in self.display_names:
            # Update internal data structures
            self.display_names.append(display_name)
            self.display_to_filename_map[display_name] = original_filename
            
            # Update the ComboBox UI
            self.cmb.addItem(display_name)
            
        # 6. Set the ComboBox to the new image's name
        index = self.cmb.findText(display_name)
        if index != -1:
            self.cmb.setCurrentIndex(index)
   
    def _render_image(self):
        """Render centered scaled image into the canvas and reposition entries/arrows."""
        # ❌ scene.clear() मत करो → इससे entries delete हो जाती हैं
        for item in self.scene.items():
            if not isinstance(item, QGraphicsProxyWidget):
                self.scene.removeItem(item)

        w, h = self.canvas.width() - 4, self.canvas.height() - 4

        if self.image:
            pix = self.image.scaled(w, h, Qt.KeepAspectRatio, Qt.SmoothTransformation)
            item = QGraphicsPixmapItem(pix)
            cx = (w - pix.width()) / 2
            cy = (h - pix.height()) / 2
            item.setPos(cx, cy)
            self.scene.addItem(item)
            self._canvas_image_item = item
        else:
            # Placeholder
            self.scene.addRect(5, 5, w - 10, h - 10, pen=QPen(QColor("#cbd5e1")))
            text_item = self.scene.addText("Upload T-Shirt Image")
            font = text_item.font()
            font.setPointSize(11)   # 👈 yaha size set karna hai (default ~9 hota hai)
            text_item.setFont(font)
            text_item.setPos(w / 2 - 100, h / 2 - 20)

        # Draw arrows + entries
        self._draw_arrows()

    def _draw_arrows(self):
        """Draw arrows and place entry boxes."""
        w, h = self.canvas.width(), self.canvas.height()
        pen = QPen(QColor("#1e40af"))
        pen.setWidth(2)

        def _draw_arrowhead(scene, pen, p1, p2):
            """Draws a V-shaped arrowhead at point p2, coming from p1."""
            ARROW_SIZE = 8
            
            # Calculate the angle of the line segment
            dx = p2.x() - p1.x()
            dy = p2.y() - p1.y()
            angle = math.atan2(dy, dx)
            
            # Calculate the points for the arrowhead triangle
            p3 = QPointF(p2.x() - ARROW_SIZE * math.cos(angle - math.pi / 6),
                        p2.y() - ARROW_SIZE * math.sin(angle - math.pi / 6))
            p4 = QPointF(p2.x() - ARROW_SIZE * math.cos(angle + math.pi / 6),
                        p2.y() - ARROW_SIZE * math.sin(angle + math.pi / 6))
            
            # Create and draw the arrowhead lines
            scene.addLine(p2.x(), p2.y(), p3.x(), p3.y(), pen).setZValue(5)
            scene.addLine(p2.x(), p2.y(), p4.x(), p4.y(), pen).setZValue(5)

        # Note: Using QPointF objects for better geometry handling
        coords = {
            "collar": [QPointF(w / 2 - 80, 0.08 * h), QPointF(w / 2, 0.08 * h)],
            "left_sleeve": [QPointF(0.20 * w, 0.30 * h), QPointF(0.33 * w, 0.30 * h)],
            "right_sleeve": [QPointF(0.80 * w, 0.40 * h), QPointF(0.72 * w, 0.40 * h)],
            "center_right": [QPointF(0.75 * w, 0.70 * h), QPointF(0.60 * w, 0.70 * h)],
            "bottom_right_label": [QPointF(0.85 * w, 0.85 * h)], 
        }
        
        entry_offsets = {
            "collar": (-80, -20),
            "left_sleeve": (-85, -20), # Adjusted for potentially wider box
            "right_sleeve": (-20, -20),  # Adjusted for potentially wider box
            "center_right": (-40, -20),
            "bottom_right_label": (-40, -20), # Adjusted for wider box and far-right position
        }
        
        LINE_Z_VALUE = 5
        BOX_Z_VALUE = 10 

        for key, pts in coords.items():
            if key != "bottom_right_label":
                # 1. Draw lines and get the last line segment
                last_segment_p1 = None
                last_segment_p2 = None
                
                for i in range(len(pts) - 1):
                    p1 = pts[i]
                    p2 = pts[i + 1]
                    
                    line = self.scene.addLine(p1.x(), p1.y(), p2.x(), p2.y(), pen)
                    line.setZValue(LINE_Z_VALUE)
                    
                    # Keep track of the last segment drawn
                    last_segment_p1 = p1
                    last_segment_p2 = p2
                
                # 2. Draw the arrowhead at the end of the last line segment
                if last_segment_p1 and last_segment_p2:
                    _draw_arrowhead(self.scene, pen, last_segment_p1, last_segment_p2)


            # 3. Create the entry box (no change needed here)
            if key not in self.entries:
                entry = QLineEdit()
                entry.setFixedWidth(120) 
                entry.setAlignment(Qt.AlignCenter)
                proxy = QGraphicsProxyWidget()
                proxy.setWidget(entry)
                self.entries[key] = proxy
                self.scene.addItem(proxy)
            
            self.entries[key].setZValue(BOX_Z_VALUE)

            # 4. Set the position of the box
            frx, fry = pts[0].x(), pts[0].y() # Get coordinates from QPointF
            dx, dy = entry_offsets.get(key, (-40, -20))
            self.entries[key].setPos(frx + dx, fry + dy)

    def create_item_selection_box(self):
        group_box = QGroupBox("Item Selection")
        group_box.setFixedWidth(1700)
    
        group_box.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Preferred)
         
    
        group_layout = QVBoxLayout()
        group_box.setLayout(group_layout)

        all_layout = QHBoxLayout()
        all_layout.setSpacing(100)

        # ---- Helper function for label + widget ----
        def add_field(label_text, widget, width=None):
            layout = QHBoxLayout()
            layout.setSpacing(15)
            layout.setContentsMargins(0, 0, 0, 0)

            lbl = QLabel(label_text)
            lbl.setFixedWidth(65)

            lbl.setStyleSheet("""
                QLabel {
                    border: none;
                    border-radius: 0px;
                    padding: 3px 6px;
                    background-color: #f9f9f9;
                }
            """)
            widget.setStyleSheet("""
                QLineEdit, QComboBox{
                    border: 1px solid black;
                    border-radius: 0;
                    padding: 3px 6px;
                    background-color: #f9f9f9;
                }
            """)
           
            layout.addWidget(lbl)

            if width:
                widget.setFixedWidth(width)
            layout.addWidget(widget)
           
            all_layout.addLayout(layout)
      
        # Fabric
        self.cloth_combo = QComboBox()
        self.cloth_combo.addItems(["Shirt", "Pant", "Kurta"])
        add_field("Fabric:", self.cloth_combo, 195)
        

        # Type
        self.type_combo = QComboBox()
        self.type_combo.addItems(["Cotton", "Silk", "Linen"])
        add_field("Type:", self.type_combo, 195)

        # Color
        self.color_combo = QLineEdit("red")
        # self.color_combo.addItems(["Red", "Blue", "Green"])
        add_field("Color:", self.color_combo, 100)

        # Size
        self.size_combo = QComboBox()
        self.size_combo.addItems(["S", "M", "L", "XL", "XXL"])
        add_field("Size:", self.size_combo, 100)

        # Quantity
        self.qty_input = QLineEdit("1")
        add_field("QTY:", self.qty_input, 60)
       
         
        #  price.unit
        self.price_input=QLineEdit("200")
        add_field("Unit:", self.price_input, 60)
        
        # Add Row Button
        self.add_button = QPushButton("+Add")
        self.add_button.setStyleSheet("background-color:#87CEFA;")
        self.add_button.setSizePolicy(QSizePolicy.Fixed, QSizePolicy.Fixed) 
        self.add_button.clicked.connect(self._add_item_row)
        all_layout.addWidget(self.add_button,0,Qt.AlignLeft)
        group_layout.addSpacing(20)

        group_layout.addLayout(all_layout)
         # --- Table aligned exactly under inputs ---
        self.items_container = QTableWidget()
        # self.items_container.setStyleSheet("QTableWidget { background-color: #f9f9f9; }")
        self.items_container.setColumnCount(7)
        # Table background aur border
        self.items_container.setStyleSheet("""
        QTableWidget {
            gridline-color: gray;  /* cells ke beech line */
            border: 0px solid black; /* table ke outer border */
        }
        QTableWidget::item {
            border: 0.5px solid gray;
            background-color: #ffffff;                                  /* cell borders */
        }
        QHeaderView::section {
            background-color:: gray; 
            border: 0.5px solid #000000;   /* header borders */
            padding: 3px;
            font-weight: normal;         /* header text bold */
            text-align: center;        /* header text center */
        }
        """)
        
        self.items_container.setHorizontalHeaderLabels(
            ["Fabric", "Type", "Color", "Size", "Qty", "Unit", "Total price"]
        )
        self.items_container.verticalHeader().setVisible(False)
        self.items_container.setFixedHeight(200)

        # Align column widths exactly like inputs above
        self.items_container.setColumnWidth(0, 300)  # Cloth
        self.items_container.setColumnWidth(1, 350)  # Type
        self.items_container.setColumnWidth(2, 200)  # Color
        self.items_container.setColumnWidth(3, 150)   # Size
        self.items_container.setColumnWidth(4, 150)   # Qty
        self.items_container.setColumnWidth(5, 200)   # Unit
        self.items_container.setColumnWidth(6, 200)   # Total

        group_layout.addWidget(self.items_container)
       
        # --- Grand Total Label ---
        self.grand_total_label = QLabel("Grand Total: 0")
        self.grand_total_label.setFixedWidth(200)
        self.grand_total_label.setStyleSheet("""
            QLabel {
                font-size: 20px;
                font-weight: bold;
                color: #000000;      
                background-color: #f9f9f9;                                                                                                       
            }
        """)
        grand_total_layout = QHBoxLayout()
        grand_total_layout.addSpacing(1200)  
        grand_total_layout.addWidget(self.grand_total_label)

        group_layout.addLayout(grand_total_layout)
        return group_box
    
        
    def _add_item_row(self):
        row = self.items_container.rowCount()
        self.items_container.insertRow(row)

        try:
            qty = int(self.qty_input.text())
            unit = float(self.price_input.text())
        except ValueError:
            print("Error: Quantity or Unit Price must be valid numbers.")
            self.items_container.removeRow(row)
            return
        printing_add_on_per_unit = self.get_total_printing_price()
        collar_add_on_per_unit = self.get_selected_collar_price()
        total = (unit + printing_add_on_per_unit + collar_add_on_per_unit) * qty

        # Cloth
        item = QTableWidgetItem(self.cloth_combo.currentText())
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 0, item)

        # Type
        item = QTableWidgetItem(self.type_combo.currentText())
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 1, item)

        # Color
        item = QTableWidgetItem(self.color_combo.text())
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 2, item)

        # Size
        item = QTableWidgetItem(self.size_combo.currentText())
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 3, item)

        # Qty
        item = QTableWidgetItem(str(qty))
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 4, item)

        # Unit
        item = QTableWidgetItem(f"{unit:.2f}")
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 5, item)

        # Total
        item = QTableWidgetItem(f"{total:.2f}")
        item.setTextAlignment(Qt.AlignCenter)
        self.items_container.setItem(row, 6, item)
        item_print = QTableWidgetItem(f"{printing_add_on_per_unit:.2f}")
        self.items_container.setItem(row, 7, item_print)       
        item_collar = QTableWidgetItem(f"{collar_add_on_per_unit:.2f}")
        self.items_container.setItem(row, 8, item_collar)
        self.items_container.setColumnHidden(7, True)
        self.items_container.setColumnHidden(8, True)
        self._update_grand_total()

    def get_total_printing_price(self):
        """Calculates the sum of prices for all selected printing options (per unit)."""
        total_print_price = 0.0
        for key, (checkbox, price_edit) in self.print_vars.items():
            if checkbox.isChecked():
                try:
                    price = float(price_edit.text())
                    total_print_price += price
                except ValueError:
                    print(f"Warning: Invalid price found for {key}.")
                    continue
        return total_print_price

    def get_selected_collar_price(self):
        """Returns the price of the currently selected collar option (per unit)."""
        collar_type = self.collar_var  # 'self', 'rib', or 'patti'
        collar_price = 0.0
        try:
            if collar_type == "self":
                collar_price = float(self.collar_price_self.text())
            elif collar_type == "rib":
                collar_price = float(self.collar_price_rib.text())
            elif collar_type == "patti":
                collar_price = float(self.collar_price_patti.text())
        except ValueError:
            print(f"Warning: Invalid price found for {collar_type} collar.")
            pass # Price remains 0.0 if conversion fails
        return collar_price

    def _calculate_item_total_price(self, unit_price, qty):        
        # Add-on prices are per unit
        printing_add_on_per_unit = self.get_total_printing_price()
        collar_add_on_per_unit = self.get_selected_collar_price()       
        add_on_total = (printing_add_on_per_unit + collar_add_on_per_unit) * qty
        base_total = unit_price * qty
        
        return base_total + add_on_total
    
    def _recalculate_all_item_totals(self):
        """Recalculates the Total Price for all rows in the table and updates the Grand Total."""
        from PyQt5.QtWidgets import QTableWidgetItem
        
        for row in range(self.items_container.rowCount()):
            unit_price_item = self.items_container.item(row, 5) 
            qty_item = self.items_container.item(row, 4) 
            if unit_price_item and qty_item:
                try:
                    unit_price = float(unit_price_item.text())
                    qty = int(qty_item.text())
                    new_total_price = self._calculate_item_total_price(unit_price, qty)
                    total_price_item = self.items_container.item(row, 6)
                    if not total_price_item:
                        total_price_item = QTableWidgetItem()
                        self.items_container.setItem(row, 6, total_price_item)
                        total_price_item.setText(f"{new_total_price:.2f}") 
                except ValueError:
                    continue 
        self._update_grand_total()
        
    def _update_grand_total(self):
        total_sum = 0
        for row in range(self.items_container.rowCount()):
            item = self.items_container.item(row, 6)  # column 6 = Total price
            if item:
                try:
                    total_sum += float(item.text())
                except ValueError:
                    pass
        self.grand_total_label.setText(f"Grand Total: {total_sum}")
    
     # --- Remark Field as separate method ---
    def add_remark_field(self):
        self.main_layout.addSpacing(10)
        remark_layout = QHBoxLayout()
        remark_layout.setSpacing(10)

        remark_label = QLabel("REMARK:")
        remark_label.setFixedWidth(100)
        remark_label.setAlignment(Qt.AlignCenter)
        remark_label.setStyleSheet("""
            QLabel {
                font-size: 16px;
                font-weight: bold;
                border: 1px solid black;
                padding: 2px 4px;
                border-radius: 4px;  /* chhota rounded corner optional */
                background-color: #f9f9f9;  /* optional background color */
            }
        """)
       
        self.remark_input = QLineEdit()
        self.remark_input.setPlaceholderText("Enter remark here...")
        self.remark_input.setFixedWidth(500)
        self.remark_input.setStyleSheet("""
            QLineEdit {
                border: 1px solid black;
                border-radius: 0px;
                padding: 4px 6px;
                background-color: #ffffff;
            }
        """)

        remark_layout.addWidget(remark_label)
        remark_layout.addWidget(self.remark_input)
        remark_layout.addStretch()

        self.main_layout.addLayout(remark_layout)
        self.main_layout.addSpacing(10)

    def create_buttons_row(self):
        buttons_layout = QHBoxLayout()
        buttons_layout.setContentsMargins(0, 0, 0, 0)
        buttons_layout.setSpacing(1)    
        
        self.save_btn = QPushButton("💾\n SAVE")
        self.undo_btn=QPushButton("↩️\n CANCEL")
        self.bill_btn = QPushButton("🧾\n GENERATE BILL")
        self.job_btn = QPushButton("⚒️\n JOB WORK")
        self.rib_btn = QPushButton("🧵\n RIB COLLAR")
        self.cut_btn = QPushButton("✂️\n CUTTING")
        self.print_btn = QPushButton("🖨\n PRINTING")
        
        self.bill_btn.setFixedWidth(170)
        self.save_btn.setFixedWidth(140)
        self.undo_btn.setFixedWidth(140)
        self.job_btn.setFixedWidth(140)
        self.rib_btn.setFixedWidth(140)
        self.cut_btn.setFixedWidth(140)
        self.print_btn.setFixedWidth(140)
       
        buttons = [
            self.save_btn,self.undo_btn, self.bill_btn, self.job_btn,
            self.rib_btn, self.cut_btn, self.print_btn
        ]

        for i, btn in enumerate(buttons):
            btn.setFixedHeight(60)
            btn.setStyleSheet(" background-color: #EEEEEE; "
           " color: #fc83a0;" 
            "font-weight: bold; "
            "padding:10px;")
            
            buttons_layout.addWidget(btn)
        
            # har button ke baad space chahiye (last ke baad nahi)
            if i < len(buttons) - 1:
                buttons_layout.addSpacing(1)  # yaha 30px gap set karo
        buttons_layout.addStretch()
        return buttons_layout
    
    def open_search_window(self):
        # नया window बनाओ
        self.search_window = QWidget()
        self.search_window.setWindowTitle("📋 Order List")
        self.search_window.setGeometry(200, 200, 500, 400)

        layout = QVBoxLayout(self.search_window)

        # Header
        header = QLabel("Select an Order from the list")
        header.setStyleSheet("font-size: 14pt; font-weight: bold;")
        layout.addWidget(header, alignment=Qt.AlignCenter)

        # Search bar
        self.search_input = QLineEdit()
        self.search_input.setPlaceholderText("🔍 Search...")
        layout.addWidget(self.search_input)

        # Order list
        self.list_widget = QListWidget()
        orders = [
            {"id": 1, "customer": "Ashwini", "total": 500},
            {"id": 2, "customer": "Rahul", "total": 700},
            {"id": 3, "customer": "Sneha", "total": 300},
        ]
        self.orders = orders
        for order in orders:
            item_text = f"Order {order['id']} | {order['customer']} | ₹{order['total']}"
            item = QListWidgetItem(item_text)
            item.setData(Qt.UserRole, order)
            self.list_widget.addItem(item)
        layout.addWidget(self.list_widget)

        # Close button
        close_btn = QPushButton("Close")
        layout.addWidget(close_btn, alignment=Qt.AlignRight)

        # Connections
        def filter_list(text):
            for i in range(self.list_widget.count()):
                item = self.list_widget.item(i)
                item.setHidden(text.lower() not in item.text().lower())

        self.search_input.textChanged.connect(filter_list)
        close_btn.clicked.connect(self.search_window.close)

        def on_item_selected(item):
            order = item.data(Qt.UserRole)
            # 👇 main form ke fields fill
            self.party_name.setText(order["customer"])
            self.total_price_input.setText(str(order["total"]))
            self.search_window.close()

        self.list_widget.itemDoubleClicked.connect(on_item_selected)

        # Show new window
        self.search_window.show()

    def open_print_dialog(self):
        order_details_string = "Example item: Shirt, Red, Size S, Qty 1" 
        dialog = PrintExportDialog(content_data=order_details_string, parent=self)
        dialog.exec_()
    
   
class PrintExportDialog(QDialog):
    def __init__(self, content_data, parent=None):
        super().__init__(parent)
        self.setWindowTitle("Print and Export Options")
        self.content_data = content_data # Data to be printed/exported
        self.setGeometry(200, 200, 350, 200)

        main_layout = QVBoxLayout(self)

        # Print Options
        print_layout = QHBoxLayout()
        self.direct_print_btn = QPushButton("🖨 Direct Print")
        self.preview_btn = QPushButton("🔍 Print Preview")
        
        print_layout.addWidget(self.direct_print_btn)
        print_layout.addWidget(self.preview_btn)
        main_layout.addLayout(print_layout)

        # Connect Signals
        self.direct_print_btn.clicked.connect(self.direct_print)
        self.preview_btn.clicked.connect(self.show_preview)
        
    def get_print_content(self):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        party_name = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"

        html_content = f"""
        <html><body>
            <h1>Order Report</h1>
            <p>Order No: {self.parent().order_number.text()}</p>
            <p>Party Name: {self.parent().party_name.text()}</p>
            <hr>
            <h2>Item Details:</h2>
            <p>{self.content_data}</p>
        </body></html>
        """
        return html_content

    def print_document(self, printer):
        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        doc.print_(printer)

    def direct_print(self):
        printer = QPrinter(QPrinter.HighResolution)
        dialog = QPrintDialog(printer, self)
        if dialog.exec_() == QDialog.Accepted:
            self.print_document(printer)

    def show_preview(self):
        """Opens the built-in Print Preview dialog."""
        printer = QPrinter(QPrinter.HighResolution)
        preview = QPrintPreviewDialog(printer, self)
        preview.paintRequested.connect(self.print_document)

        export_btn = QPushButton("🔽 Export Options")
        export_btn.setToolTip("Export to PDF/Excel/Image/Word/PPT")
        export_btn.clicked.connect(self.show_export_menu_from_preview)
        preview.layout().addWidget(export_btn) 

        share_btn = QPushButton("📱 Share via WhatsApp")
        share_btn.setToolTip("Share Order as PDF via WhatsApp")
        share_btn.clicked.connect(self.show_whatsapp_share_menu_from_preview)
        preview.layout().addWidget(share_btn)

        preview.exec_()

    def show_export_menu_from_preview(self):

        menu = QMenu(self)
        # PDF Export
        pdf_action = menu.addAction("Export to PDF (*.pdf)")
        pdf_action.triggered.connect(lambda: self._perform_pdf_save(None, show_msg=True))
        # Image Export
        image_action = menu.addAction("Export to Image (*.png, *.jpg)")
        image_action.triggered.connect(lambda: self._perform_image_save(None, show_msg=True)) 
        menu.addSeparator()
        # Excel Export
        excel_action = menu.addAction("Export to Excel (*.xlsx)")
        excel_action.triggered.connect(lambda: self._perform_excel_save(None, show_msg=True))
        # Word Export
        word_action = menu.addAction("Export to Word (*.docx)")
        word_action.triggered.connect(lambda: self._perform_word_ppt_save(None, 'word', show_msg=True))
        # PowerPoint Export
        ppt_action = menu.addAction("Export to PowerPoint (*.pptx)")
        ppt_action.triggered.connect(lambda: self._perform_word_ppt_save(None, 'ppt', show_msg=True))
        menu.exec_(QCursor.pos())

    def _perform_pdf_save(self, fileName=None, show_msg=False):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"  
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save PDF for Sharing", f"Order_{order_no}.pdf", "PDF Files (*.pdf)"
            )
            if not fileName: return None
            show_msg = True 
        printer = QPrinter(QPrinter.HighResolution)
        printer.setOutputFormat(QPrinter.PdfFormat)
        printer.setOutputFileName(fileName)
        self.print_document(printer)        
        if show_msg:
            QMessageBox.information(self, "Success", f"PDF saved to:\n{fileName}")
        return fileName
    
    def _perform_excel_save(self, fileName=None, show_msg=False):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"    
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save Excel for Sharing", f"Order_{order_no}.xlsx", "Excel Files (*.xlsx)"
            )
            if not fileName: return None
            show_msg = True
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws['A1'] = "Order No:"
            ws['B1'] = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
            ws['A2'] = "Party Name:"
            ws['B2'] = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"
            ws['A4'] = "Item Details (Raw Content):"
            ws['A5'] = self.content_data 
            wb.save(fileName)
            if show_msg:
                QMessageBox.information(self, "Success", f"Excel file saved to:\n{fileName}")
            return fileName

        except ImportError:
            QMessageBox.critical(self, "Error", "The 'openpyxl' library is required for Excel export. Please install it.")
            return None
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during Excel export: {e}")
            return None

    def _perform_image_save(self, fileName=None, show_msg=False):

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save Image for Sharing", f"Order_{order_no}.png", "Image Files (*.png);;JPEG Files (*.jpg)"
            )
            if not fileName: return None
            show_msg = True        

        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        image_size = doc.size().toSize()
        if image_size.isEmpty(): image_size = QSize(800, 1000)           
        image = QPixmap(image_size)
        image.fill(self.palette().window().color())
        painter = QPainter(image)
        doc.drawContents(painter, QRectF(image.rect())) 
        painter.end()
        image.save(fileName)
        if show_msg:
            QMessageBox.information(self, "Success", f"Image saved to:\n{fileName}")
        return fileName

    def _perform_word_ppt_save(self, fileName=None, file_type='word', show_msg=False):
        
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        party_name = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"
        
        if file_type == 'word':
            filter_str, default_ext, title_verb = "Word Documents (*.docx)", ".docx", "Document"
        else: # ppt
            filter_str, default_ext, title_verb = "PowerPoint Presentations (*.pptx)", ".pptx", "Presentation"
        
        if fileName is None:
            title = f"Save {file_type.title()} for Sharing"
            fileName, _ = QFileDialog.getSaveFileName(
                self, title, f"Order_{order_no}{default_ext}", filter_str
            )
            if not fileName: return None
            show_msg = True

        try:
            if file_type == 'word':
                from docx import Document    
                doc = Document()
                doc.add_heading('Order Report', 0)
                doc.add_paragraph(f'Order No: {order_no}')
                doc.add_paragraph(f'Party Name: {party_name}')
                doc.add_heading('Item Details:', level=2)
                doc.add_paragraph(self.content_data) 
                
                doc.save(fileName)

            elif file_type == 'ppt':
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = "Order Report"
                
                body = slide.placeholders[1]
                body.text = f"Order No: {order_no}\nParty Name: {party_name}\n\nItem Details:\n{self.content_data}"               
                prs.save(fileName)
            if show_msg:
                QMessageBox.information(self, "Success", f"{title_verb} saved to:\n{fileName}")
            return fileName
        except ImportError:
            QMessageBox.critical(self, "Error", f"The required library for {file_type.title()} export is missing. Please install 'python-docx' or 'python-pptx'.")
            return None
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during {file_type.title()} export: {e}")
            return None
        
    def share_via_whatsapp(self, file_path=None):
        """
        Opens WhatsApp Web/Desktop with a pre-filled message.
        Requires manual attachment of the file saved at file_path.
        """
        if file_path is None:
            QMessageBox.warning(self, "Error", "Please export the file first.")
            return

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        
        # Determine file type for the message
        ext = file_path.split('.')[-1].upper()
        
        message = f"Please find the Order Report (Order No: {order_no}) in {ext} format."
        
        encoded_message = QUrl.toPercentEncoding(message).data().decode()
        # NOTE: Using web.whatsapp.com for web/desktop sharing
        url = f"https://web.whatsapp.com/send?text={encoded_message}" 
        
        # Open the browser
        webbrowser.open(url)
        QMessageBox.information(self, "Manual Step Required", 
                                f"Your browser has opened WhatsApp. Please **manually attach** the saved file:\n\n{file_path}")

    def share_office_format_via_whatsapp(self, file_type):
        """Saves and shares Excel or Word/PPT via WhatsApp."""
        file_path = None
        
        if file_type == 'excel':
            file_path = self._perform_excel_save(None) # Pass None to prompt file dialog
        elif file_type == 'word':
            file_path = self._perform_word_ppt_save(None, 'word')
        
        if file_path:
            self.share_via_whatsapp(file_path)
            
    def share_via_whatsapp(self, file_path=None):
        """Opens WhatsApp link, prompting user to manually attach the file."""
        if file_path is None:
            QMessageBox.warning(self, "Error", "Please export the file first.")
            return

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        ext = file_path.split('.')[-1].upper()
        message = f"Please find the Order Report (Order No: {order_no}) in {ext} format."
        encoded_message = QUrl.toPercentEncoding(message).data().decode()
        url = f"https://web.whatsapp.com/send?text={encoded_message}" 
        # 1. Open the WhatsApp link (opens browser/desktop app)
        webbrowser.open(url)

        # 2. Open the file's containing folder (Cross-platform way to show the file)
        try:
            if sys.platform == 'win32':
                # Windows: Opens Explorer and selects the file
                os.startfile(os.path.dirname(file_path)) 
            elif sys.platform == 'darwin':
                # macOS: Opens Finder and selects the file
                os.system(f'open -R "{file_path}"')
            else:
                # Linux: Opens the file's directory
                os.system(f'xdg-open "{os.path.dirname(file_path)}"')
                
            QMessageBox.information(
                self, 
                "Action Required 🚨", 
                f"1. WhatsApp has opened in your browser/desktop app.\n"
                f"2. The file folder has also opened.\n\n"
                f"Please **drag and drop** the file onto the chat window to share it!"
            )
        except Exception as e:
            QMessageBox.warning(
                self, 
                "Action Required 🚨", 
                f"WhatsApp link opened. Could not open file explorer automatically.\n\n"
                f"Please manually navigate to and attach the saved file:\n\n{file_path}"
            )

    def show_whatsapp_share_menu_from_preview(self):
        """Shows format options for sharing, using the reusable save functions."""
        menu = QMenu(self)
        
        pdf_action = menu.addAction("Share as PDF")
        pdf_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_pdf_save))
        
        image_action = menu.addAction("Share as Image (PNG/JPG)")
        image_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_image_save))
        
        menu.addSeparator()
        excel_action = menu.addAction("Share as Excel (.xlsx)")
        excel_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_excel_save))
        
        word_action = menu.addAction("Share as Word (.docx)")
        word_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_word_ppt_save, 'word'))

        menu.exec_(QCursor.pos())

    def share_file_and_whatsapp(self, save_func, file_type=None):
        
        if file_type:
            file_path = save_func(None, file_type)
        else:
            file_path = save_func(None)
            
        if file_path:
            self.share_via_whatsapp(file_path)

if __name__ == "__main__":
    app = QApplication(sys.argv)

    app.setStyleSheet("""
        QWidget {
            font-size: 11pt;
            font-family: Arial;
        }
    """)
    window = OrderForm()
    window.show()
    sys.exit(app.exec_())