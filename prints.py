import sys
import os
import webbrowser
from PyQt5.QtWidgets import (QDialog, QVBoxLayout, QHBoxLayout, QPushButton, QMessageBox, QFileDialog, QMenu, QApplication)
from PyQt5.QtPrintSupport import QPrinter, QPrintDialog, QPrintPreviewDialog
from PyQt5.QtGui import QTextDocument, QCursor, QPixmap, QPainter
from PyQt5.QtCore import QUrl, QSize, QRectF, QDate
from openpyxl import Workbook
from docx import Document
from pptx import Presentation
from pptx.util import Inches
import re

class ExportShareMixin:

    def show_export_menu(self):
        """Shows the format options for local file saving."""
        menu = QMenu(self)

        pdf_action = menu.addAction("Save to PDF (*.pdf)")
        # Note: Set show_msg=True to confirm save, as this is a direct save action
        pdf_action.triggered.connect(lambda: self._perform_pdf_save(None, show_msg=True))

        image_action = menu.addAction("Save to Image (*.png, *.jpg)")
        image_action.triggered.connect(lambda: self._perform_image_save(None, show_msg=True)) 
        menu.addSeparator()

        excel_action = menu.addAction("Save to Excel (*.xlsx)")
        excel_action.triggered.connect(lambda: self._perform_excel_save(None, show_msg=True))

        word_action = menu.addAction("Save to Word (*.docx)")
        word_action.triggered.connect(lambda: self._perform_word_ppt_save(None, 'word', show_msg=True))

        ppt_action = menu.addAction("Save to PowerPoint (*.pptx)")
        ppt_action.triggered.connect(lambda: self._perform_word_ppt_save(None, 'ppt', show_msg=True))
        
        menu.exec_(QCursor.pos())

    def _perform_pdf_save(self, fileName=None, show_msg=False):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp" 
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save PDF for Sharing", f"Order_{order_no}.pdf", "PDF Files (*.pdf)"
            )
            if not fileName: return None
            # If the user chose a name via dialog, we want to show a success message
            show_msg = True 
        
        printer = QPrinter(QPrinter.HighResolution)
        printer.setOutputFormat(QPrinter.PdfFormat)
        printer.setOutputFileName(fileName)
        self.print_document(printer) # <-- REQUIRES print_document METHOD
        
        if show_msg:
            QMessageBox.information(self, "Success", f"PDF saved to:\n{fileName}")
        return fileName

    def _perform_excel_save(self, fileName=None, show_msg=False):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp" 
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save Excel for Sharing", f"Order_{order_no}.xlsx", "Excel Files (*.xlsx)"
            )
            if not fileName: return None
            show_msg = True
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            # Simplified Excel content extraction (adapt as needed)
            ws['A1'] = "Order No:"
            ws['B1'] = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
            ws['A2'] = "Party Name:"
            ws['B2'] = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"
            ws['A4'] = "Item Details (Raw Content):"
            # NOTE: self.content_data is accessed directly from the inheriting class
            ws['A5'] = self.content_data 
            wb.save(fileName)
            if show_msg:
                QMessageBox.information(self, "Success", f"Excel file saved to:\n{fileName}")
            return fileName

        except ImportError:
            QMessageBox.critical(self, "Error", "The 'openpyxl' library is required for Excel export. Please install it.")
            return None
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during Excel export: {e}")
            return None

    def _perform_image_save(self, fileName=None, show_msg=False):

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save Image for Sharing", f"Order_{order_no}.png", "Image Files (*.png);;JPEG Files (*.jpg)"
            )
            if not fileName: return None
            show_msg = True        

        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        image_size = doc.size().toSize()
        if image_size.isEmpty(): image_size = QSize(800, 1000)           
        image = QPixmap(image_size)
        image.fill(self.palette().window().color())
        painter = QPainter(image)
        doc.drawContents(painter, QRectF(image.rect())) 
        painter.end()
        image.save(fileName)
        if show_msg:
            QMessageBox.information(self, "Success", f"Image saved to:\n{fileName}")
        return fileName

    def _perform_word_ppt_save(self, fileName=None, file_type='word', show_msg=False):
        
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        party_name = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"
        
        if file_type == 'word':
            filter_str, default_ext, title_verb = "Word Documents (*.docx)", ".docx", "Document"
        else: # ppt
            filter_str, default_ext, title_verb = "PowerPoint Presentations (*.pptx)", ".pptx", "Presentation"
        
        if fileName is None:
            title = f"Save {file_type.title()} for Sharing"
            fileName, _ = QFileDialog.getSaveFileName(
                self, title, f"Order_{order_no}{default_ext}", filter_str
            )
            if not fileName: return None
            show_msg = True

        try:
            if file_type == 'word':
                from docx import Document    
                doc = Document()
                doc.add_heading('Order Report', 0)
                doc.add_paragraph(f'Order No: {order_no}')
                doc.add_paragraph(f'Party Name: {party_name}')
                doc.add_heading('Item Details:', level=2)
                doc.add_paragraph(self.content_data) 
                
                doc.save(fileName)

            elif file_type == 'ppt':
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = "Order Report"
                
                body = slide.placeholders[1]
                body.text = f"Order No: {order_no}\nParty Name: {party_name}\n\nItem Details:\n{self.content_data}"               
                prs.save(fileName)
            if show_msg:
                QMessageBox.information(self, "Success", f"{title_verb} saved to:\n{fileName}")
            return fileName
        except ImportError:
            QMessageBox.critical(self, "Error", f"The required library for {file_type.title()} export is missing. Please install 'python-docx' or 'python-pptx'.")
            return None
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during {file_type.title()} export: {e}")
            return None

    # The correct, detailed share_via_whatsapp function
    def share_via_whatsapp(self, file_path=None):
        """Opens WhatsApp link, prompting user to manually attach the file."""
        if file_path is None:
            QMessageBox.warning(self, "Error", "Please export the file first.")
            return

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        ext = file_path.split('.')[-1].upper()
        message = f"Please find the Order Report (Order No: {order_no}) in {ext} format."
        encoded_message = QUrl.toPercentEncoding(message).data().decode()
        url = f"https://web.whatsapp.com/send?text={encoded_message}" 
        
        # 1. Open the WhatsApp link
        webbrowser.open(url)

        # 2. Open the file's containing folder
        try:
            if sys.platform == 'win32':
                os.startfile(os.path.dirname(file_path)) 
            elif sys.platform == 'darwin':
                os.system(f'open -R "{file_path}"')
            else:
                os.system(f'xdg-open "{os.path.dirname(file_path)}"')
                
            QMessageBox.information(
                self, 
                "Action Required 🚨", 
                f"1. WhatsApp has opened in your browser/desktop app.\n"
                f"2. The file folder has also opened.\n\n"
                f"Please **drag and drop** the file onto the chat window to share it!"
            )
        except Exception as e:
            QMessageBox.warning(
                self, 
                "Action Required 🚨", 
                f"WhatsApp link opened. Could not open file explorer automatically.\n\n"
                f"Please manually navigate to and attach the saved file:\n\n{file_path}"
            )

    def show_whatsapp_share_menu(self):
        """Shows format options for sharing, using the reusable save functions."""
        menu = QMenu(self)
        
        pdf_action = menu.addAction("Share as PDF")
        pdf_action.triggered.connect(lambda: self._share_file_and_whatsapp(self._perform_pdf_save))
        
        image_action = menu.addAction("Share as Image (PNG/JPG)")
        image_action.triggered.connect(lambda: self._share_file_and_whatsapp(self._perform_image_save))
        
        menu.addSeparator()
        excel_action = menu.addAction("Share as Excel (.xlsx)")
        excel_action.triggered.connect(lambda: self._share_file_and_whatsapp(self._perform_excel_save))
        
        word_action = menu.addAction("Share as Word (.docx)")
        word_action.triggered.connect(lambda: self._share_file_and_whatsapp(self._perform_word_ppt_save, 'word'))
        
        # NOTE: ppt_action is included here for completeness, though often not shared via WA
        # ppt_action = menu.addAction("Share as PowerPoint (.pptx)")
        # ppt_action.triggered.connect(lambda: self._share_file_and_whatsapp(self._perform_word_ppt_save, 'ppt'))


        menu.exec_(QCursor.pos())

    def _share_file_and_whatsapp(self, save_func, file_type=None):
        """Internal helper to call the save function and then the share function."""
        file_path = None
        
        # Fix the argument passing issue from the previous response
        if save_func == self._perform_word_ppt_save:
            # Need to pass file_type and explicitly set show_msg=False
            file_path = save_func(None, file_type, show_msg=False) 
        else:
            # All other save functions accept (fileName, show_msg)
            file_path = save_func(None, show_msg=False) 
            
        if file_path:
            self.share_via_whatsapp(file_path)

class PrintExportDialog(QDialog):
    def __init__(self, parent, content_data, document_type="ORDER", **kwargs):
        super().__init__(parent, **kwargs)
        self.setWindowTitle("Print and Export Options")
        self.content_data = content_data 
        self.document_type = document_type
        self.setGeometry(200, 200, 350, 200)

        main_layout = QVBoxLayout(self)

        # Print Options
        print_layout = QHBoxLayout()
        self.direct_print_btn = QPushButton("🖨 Direct Print")
        self.preview_btn = QPushButton("🔍 Print Preview")
        
        print_layout.addWidget(self.direct_print_btn)
        print_layout.addWidget(self.preview_btn)
        main_layout.addLayout(print_layout)

        # Connect Signals
        self.direct_print_btn.clicked.connect(self.direct_print)
        self.preview_btn.clicked.connect(self.show_preview)

    def _get_parent_text(self, attribute_name, default="N/A"):
        parent = self.parent()
        if hasattr(parent, attribute_name) and getattr(parent, attribute_name) and hasattr(getattr(parent, attribute_name), 'text'):
            return getattr(parent, attribute_name).text()
        return default

    def _get_parent_checkbox_state(self, checkbox_attr_name, price_attr_name):
        parent = self.parent()
        if hasattr(parent, checkbox_attr_name):
            widget = getattr(parent, checkbox_attr_name)
            if hasattr(widget, 'isChecked') and widget.isChecked():
                if price_attr_name is None: 
                    return widget.text().strip()
                price = self._get_parent_text(price_attr_name, "0")
                option_text = widget.text().strip()
                return f"{option_text} (Price: {price})"
        return None
    
    def _get_parent_printing_options(self):
        parent = self.parent()
        printing_options = []
        
        keys = ['front', 'back', 'patch', 'embroidery', 'Dtf', 'Front sablimation', 'Back sablimation']
        
        if hasattr(parent, 'print_vars'):
            print_vars = getattr(parent, 'print_vars')
            for key in keys:
                if key in print_vars:
                    checkbox, price_edit = print_vars[key]
                    
                    if checkbox.isChecked():
                        price = price_edit.text() if hasattr(price_edit, 'text') else "0"
                        option_text = checkbox.text().strip()
                        printing_options.append(f"<li>{option_text} (Price: {price} INR)</li>")

        return "".join(printing_options) or "<li>None Selected</li>"
    
    def _get_parent_track_options(self):
        parent = self.parent()
        track_options_list = []
        
        if hasattr(parent, 'track_vars') and hasattr(parent, 'track_extra_vars'):
            track_vars = getattr(parent, 'track_vars')
            track_extra_vars = getattr(parent, 'track_extra_vars')
            
            # NOTE: This key list must match the keys used in _build_options_panel
            keys = ['Dori', '1 Piping', '2 Piping', 'Other'] 
            
            for key in keys:
                if key in track_vars:
                    checkbox, price_edit = track_vars[key]
                    extra_edit = track_extra_vars.get(key)
                    
                    if checkbox.isChecked():
                        price = price_edit.text() if hasattr(price_edit, 'text') else "0"
                        extra_detail = extra_edit.text() if extra_edit and hasattr(extra_edit, 'text') else ""
                        
                        option_text = checkbox.text().strip()
                        
                        detail_info = f" ({extra_detail})" if extra_detail else ""
                        
                        track_options_list.append(f"<li>{option_text}{detail_info} (Price: {price} INR)</li>")

        return "".join(track_options_list) or "<li>None Selected</li>"
    
    # --- NEW METHOD FOR TAX CALCULATION ---
    def _get_tax_info(self):
        main_window = self.parent()
        
        try:
            total_items_price = main_window._total_items_sum
            tax_percent = main_window._tax_percentage
            tax_amount = main_window._tax_amount
            grand_total = main_window._grand_total # The final, calculated number
        except AttributeError:
            total_items_price = 0.0
            tax_percent = 0.0
            tax_amount = 0.0
            grand_total = 0.0

        tax_apply_text = self._get_parent_text('tax_apply_combo')
        display_grand_total = f"{grand_total:.2f}"

        tax_summary_html = f"""
        <table style="width: 100%; text-align: right; border-collapse: collapse; margin-top: 10px;">
            <tr>
                <td style="padding: 3px 0; border-top: 1px solid #ddd;"><b>Total Items Price:</b></td>
                <td style="padding: 3px 0; border-top: 1px solid #ddd;">₹ {total_items_price:.2f}</td>
            </tr>
            <tr>
                <td style="padding: 3px 0;"><b>Tax (GST) @ {tax_percent:.1f}%:</b></td>
                <td style="padding: 3px 0;">₹ {tax_amount:.2f}</td>
            </tr>
            <tr>
                <td style="font-size: 14pt; color: #d9534f; padding: 5px 0; border-top: 2px solid #333; border-bottom: 2px solid #333;"><b>GRAND TOTAL:</b></td>
                <td style="font-size: 14pt; color: #d9534f; padding: 5px 0; border-top: 2px solid #333; border-bottom: 2px solid #333;">₹ {display_grand_total}</td>
            </tr>
        </table>
        """
        return tax_summary_html, display_grand_total
    
    def _clean_table_html(self, html_content):
        """
        A dedicated helper to remove all table columns (<th>/<td>) that appear 
        AFTER the 'Total Price' column.
        """
        if not html_content:
            return ""

        # Pattern 1: Target the Header Row (assumed to be the first <tr>...</tr>)
        def clean_header(match):
            header_row = match.group(0)
            # Find the position right after "Total Price</th>"
            # Use re.escape to handle potential special chars in "Total Price" if they were present
            pattern = r'(Total Price<\/th>)'
            
            match_end = re.search(pattern, header_row, re.IGNORECASE)
            if match_end:
                # Keep everything up to the end of "Total Price</th>", then close the row
                return header_row[:match_end.end()] + '</tr>'
            return header_row # Return original if Total Price not found
        
        # Pattern 2: Target the Data Rows (all <tr>...</tr> after the header)
        def clean_data_row(row_match):
            data_row = row_match.group(0)
            
            # Count <td>...</td> pairs to find the one for "Total Price" (it's the 7th column in the image)
            # Fabric (1), Type (2), Color (3), Size (4), Qty (5), Unit Price (6), Total Price (7)
            td_matches = list(re.finditer(r'<td.*?>.*?<\/td>', data_row, re.DOTALL | re.IGNORECASE))
            
            if len(td_matches) >= 7:
                # The 7th match is the Total Price column (index 6 since it's 0-based)
                total_price_td_end = td_matches[6].end()
                
                # Now, find the closing </tr> tag
                tr_end_match = re.search(r'<\/tr>', data_row)
                
                if tr_end_match:
                    tr_end_pos = tr_end_match.start()
                    
                    # Ensure the total_price_td_end is before the </tr>
                    if total_price_td_end < tr_end_pos:
                        # Reconstruct the row: everything before the end of 7th <td> + the final </tr>
                        return data_row[:total_price_td_end] + data_row[tr_end_pos:]
            
            return data_row # Return original if criteria not met

        # Split content into rows based on <tr> and </tr>
        rows = re.split(r'(<tr>.*?<\/tr>)', html_content, flags=re.DOTALL | re.IGNORECASE)
        
        cleaned_rows = []
        if rows:
            # The first non-empty row segment is usually the header row
            for i, row in enumerate(rows):
                if not row.strip():
                    continue

                if i == 1: # The first actual <tr>...</tr> block (the header)
                    cleaned_rows.append(clean_header(re.search(r'<tr>.*?<\/tr>', row, re.DOTALL | re.IGNORECASE)))
                
                elif '<tr>' in row and '</tr>' in row: # Data rows
                    cleaned_rows.append(clean_data_row(re.search(r'<tr>.*?<\/tr>', row, re.DOTALL | re.IGNORECASE)))
                
                else: # Any text before/after the table (like wrapper divs)
                    cleaned_rows.append(row)

        return "".join(cleaned_rows).replace('None', '') # Final cleanup and reassembly
    
    def get_print_content(self):

        order_no = self._get_parent_text('order_number')
        party_name = self._get_parent_text('party_name')
        order_date = self._get_parent_text('order_date') 
        delivery_date = self._get_parent_text('delivery_date')
        gst_no = self._get_parent_text('gst_no') 
        address = self._get_parent_text('address') 
        raw_grand_total_text = self._get_parent_text('grand_total_label')
        if ':' in raw_grand_total_text:
            grand_total = raw_grand_total_text.split(':')[-1].strip()
        else:
            grand_total = "0.00" 
        remarks = self._get_parent_text('remark_input') 
        school_name = self._get_parent_text('school_name') 
        barcode = self._get_parent_text('barcode')
        advance_paid = self._get_parent_text('advance_paid')

        parent = self.parent()
        image_base64_uri = ""
        if hasattr(parent, '_capture_canvas_as_base64'):
            image_base64_uri = parent._capture_canvas_as_base64()

        collar_options = [
            self._get_parent_checkbox_state('rb_self', 'collar_price_self'), 
            self._get_parent_checkbox_state('rb_rib', 'collar_price_rib'),
            self._get_parent_checkbox_state('rb_patti', 'collar_price_patti')
        ]
        collar_options_list = "".join(f"<li>{opt}</li>" for opt in collar_options if opt) or "<li>None Selected</li>"

        button_options = []
        button_option_map = {
            'BUTTON': 'rb_button',
            'PLAIN': 'rb_plain',
            'BOX': 'rb_box',
            'V+': 'rb_vplus'
        }
        
        for text, attr_name in button_option_map.items():
            if hasattr(parent, attr_name):
                widget = getattr(parent, attr_name)
                if hasattr(widget, 'isChecked') and widget.isChecked():
                    button_options.append(f"{text}")
                    
        button_options_list = "".join(f"<li>{opt}</li>" for opt in button_options) or "<li>None Selected (Default)</li>"

        printing_options_list = self._get_parent_printing_options()

        track_pant_options_list = self._get_parent_track_options()

        if self.document_type == "QUOTATION":
            report_title = "Quotation / Estimate"   
        else: 
            report_title = "Order Summary"

        cleaned_content_data = self._clean_table_html(self.content_data)
    
        html_content = f"""
        <html>
        <head>
            <style>
                body {{ font-family: 'Arial', sans-serif; font-size: 10pt; }}
                h1 {{ text-align: center; margin-bottom: 5px; color: #333; }}
                hr {{ border: 0.5px solid #ccc; }}
                .header-table {{ width: 100%; border-collapse: collapse; margin-bottom: 10px; }}
                .header-table td {{ padding: 3px 5px; vertical-align: top; }}
                .section-header {{ background-color: #f0f0f0; padding: 5px; margin-top: 15px; margin-bottom: 5px; border-left: 5px solid #007bff; }}
                
                .main-content-footer {{display: flex; flex-direction: row; justify-content: space-between; align-items: flex-start; width: 100%;}}
                
                .item-details-section {{flex-grow: 1; flex-basis: 60%; min-width: 0%; padding-right: 15px; box-sizing: border-box; overflow: hidden;}}

                .item-table-container {{ overflow-x: auto; }}
                .item-table {{ width: 100%; border-collapse: collapse; margin-top: 10px; }}
                .item-table th, .item-table td {{ border: 1px solid #ddd; padding: 6px; text-align: left; }}
                .item-table th {{ background-color: #e9ecef; }}

                .summary {{flex-basis: 35%; flex-basis: 35%; width: 35%; padding: 10px; border: 2px solid #333; text-align: right; margin-top: 15px;}}
                
                .options-table td {{ font-size: 9pt; }}
                .options-cell {{width: 55%; 
                    vertical-align: top; 
                    border-left: 1px solid #ddd; 
                    padding-left: 15px;
                    box-sizing: border-box; 
                    word-wrap: break-word;
                }}                
                .image-cell {{
                    width: 45%; 
                    vertical-align: top; 
                    padding-right: 15px; 
                    text-align: center;
                }}
                .product-image-preview {{ max-width: 100%; max-height: 250px; width: auto; height: auto; border: 1px solid #ccc; object-fit: contain;}}
                .options-table ul {{margin: 0 0 5px 0 !important; padding-left: 10px !important;}}
            </style>
        </head>
        <body>
            <h1>Order Summary</h1>
            <hr>

            <table class="header-table">
                <tr>
                    <td width="33%"><b>Order No:</b> {order_no}</td>
                    <td width="33%"><b>Order Date:</b> {order_date}</td>
                    <td width="34%"><b>Delivery Date:</b> {delivery_date}</td>
                </tr>
                <tr>
                    <td><b>Party Name:</b> {party_name}</td>
                    <td><b>GST No:</b> {gst_no}</td>
                </tr>
                <tr>
                    <td colspan="3"><b>School Name:</b> {school_name}</td>
                </tr>
                <tr>
                    <td colspan="3"><b>Address:</b> {address}</td>
                </tr>
                <tr>
                    
                    <td><b>Barcode:</b> {barcode}</td>
                    <td><b>Advance Paid:</b> ₹ {advance_paid}</td>
                </tr>
            </table>

            <h2 class="section-header">Product Design & Customization</h2>
            <table style="width: 100%; border-collapse: collapse; margin-bottom: 20px;">
                <tr>
                    <td class="image-cell">
                        <div style="max-width: 100%; margin: 0 auto; display: flex; align-items: center; justify-content: center;">
                            {f'<img src="{image_base64_uri}" class="product-image-preview" alt="Product Design"/>' if image_base64_uri else '<p>No Product Image Selected.</p>'}
                        </div>
                    </td>

                    <table style="width: 100%; border-collapse: collapse;">
                        <tr>
                            <td style="width: 50%; vertical-align: top; padding-right: 10px;">
                                <h3 style="margin-top: 0; margin-bottom: 5px; font-size: 11pt; color: #007bff;">Printing Options</h3>
                                <ul style="list-style-type: disc; padding-left: 20px; margin: 0 0 10px 0; font-size: 10pt;">
                                    {printing_options_list.replace('list-style-type: none;','')}
                                </ul>
            
                                <h3 style="margin-top: 0; margin-bottom: 5px; font-size: 11pt; color: #007bff;">Collar Options</h3>
                                <ul style="list-style-type: disc; padding-left: 20px; margin: 0 0 10px 0; font-size: 10pt;">
                                    {collar_options_list.replace('list-style-type: none;','')}
                                </ul>
                                
                                <h3 style="margin-top: 5px; margin-bottom: 5px; font-size: 11pt; color: #007bff;">Button Options</h3>
                                <ul style="list-style-type: disc; padding-left: 20px; margin: 0 0 0 0; font-size: 10pt;">
                                    {button_options_list.replace('list-style-type: none;','')}
                                </ul>

                                <h3 style="margin-top: 5px; margin-bottom: 5px; font-size: 11pt; color: #007bff;">Track Pant Options</h3>
                                <ul style="list-style-type: disc; padding-left: 20px; margin: 0 0 0 0; font-size: 10pt;">
                                    {track_pant_options_list}
                                </ul>
                            </td>
                        </tr>
                    </table>
                </tr>
            </table>

            <h2 class="section-header">Item Details</h2>
            
            <div class="main-content-footer">
                
                <div class="item-details-section">
                    <div class="item-table-container">
                        {cleaned_content_data}
                    </div>
                </div>

                <div class="summary">
                    {self._get_tax_info()[0]}
                </div>
            </div>
            
            <div style="clear: both; margin-top: 10px;">
                <h2 class="section-header">Remark</h2>
                <p>{remarks if remarks != "N/A" else "No special remarks."}</p>
            </div>
            
            <div style="margin-top: 50px; text-align: center; font-size: 8pt; color: #777;">
                <p>Signature (Seller)</p>
            </div>

        </body>
        </html>
        """
        return html_content
   
    def print_document(self, printer):
        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        doc.print_(printer)

    def direct_print(self):
        printer = QPrinter(QPrinter.HighResolution)
        dialog = QPrintDialog(printer, self)
        if dialog.exec_() == QDialog.Accepted:
            self.print_document(printer)

    def show_preview(self):
        printer = QPrinter(QPrinter.HighResolution)
        preview = QPrintPreviewDialog(printer, self)
        preview.paintRequested.connect(self.print_document)

        export_btn = QPushButton("🔽 Save Options")
        export_btn.setToolTip("Save to PDF/Excel/Image/Word/PPT")
        export_btn.clicked.connect(self.show_export_menu_from_preview)
        preview.layout().addWidget(export_btn) 

        share_btn = QPushButton("📱 Share via WhatsApp")
        share_btn.setToolTip("Share Order as PDF via WhatsApp")
        share_btn.clicked.connect(self.show_whatsapp_share_menu_from_preview)
        preview.layout().addWidget(share_btn)

        preview.exec_()

    def show_export_menu_from_preview(self):

        menu = QMenu(self)

        pdf_action = menu.addAction("Save to PDF (*.pdf)")
        pdf_action.triggered.connect(lambda: self._perform_pdf_save(None, show_msg=True))

        image_action = menu.addAction("Save to Image (*.png, *.jpg)")
        image_action.triggered.connect(lambda: self._perform_image_save(None, show_msg=True)) 
        menu.addSeparator()

        excel_action = menu.addAction("Save to Excel (*.xlsx)")
        excel_action.triggered.connect(lambda: self._perform_excel_save(None, show_msg=True))

        word_action = menu.addAction("Save to Word (*.docx)")
        word_action.triggered.connect(lambda: self._perform_word_ppt_save(None, 'word', show_msg=True))

        ppt_action = menu.addAction("Save to PowerPoint (*.pptx)")
        ppt_action.triggered.connect(lambda: self._perform_word_ppt_save(None, 'ppt', show_msg=True))
        menu.exec_(QCursor.pos())

    def _perform_pdf_save(self, fileName=None, show_msg=False):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"  
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save PDF for Sharing", f"Order_{order_no}.pdf", "PDF Files (*.pdf)"
            )
            if not fileName: return None
            show_msg = True 
        printer = QPrinter(QPrinter.HighResolution)
        printer.setOutputFormat(QPrinter.PdfFormat)
        printer.setOutputFileName(fileName)
        self.print_document(printer)        
        if show_msg:
            QMessageBox.information(self, "Success", f"PDF saved to:\n{fileName}")
        return fileName
    
    def _perform_excel_save(self, fileName=None, show_msg=False):
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"    
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save Excel for Sharing", f"Order_{order_no}.xlsx", "Excel Files (*.xlsx)"
            )
            if not fileName: return None
            show_msg = True
        try:
            from openpyxl import Workbook
            wb = Workbook()
            ws = wb.active
            ws['A1'] = "Order No:"
            ws['B1'] = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
            ws['A2'] = "Party Name:"
            ws['B2'] = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"
            ws['A4'] = "Item Details (Raw Content):"
            ws['A5'] = self.content_data 
            wb.save(fileName)
            if show_msg:
                QMessageBox.information(self, "Success", f"Excel file saved to:\n{fileName}")
            return fileName

        except ImportError:
            QMessageBox.critical(self, "Error", "The 'openpyxl' library is required for Excel export. Please install it.")
            return None
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during Excel export: {e}")
            return None

    def _perform_image_save(self, fileName=None, show_msg=False):

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "temp"
        if fileName is None:
            fileName, _ = QFileDialog.getSaveFileName(
                self, "Save Image for Sharing", f"Order_{order_no}.png", "Image Files (*.png);;JPEG Files (*.jpg)"
            )
            if not fileName: return None
            show_msg = True        

        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        image_size = doc.size().toSize()
        if image_size.isEmpty(): image_size = QSize(800, 1000)           
        image = QPixmap(image_size)
        image.fill(self.palette().window().color())
        painter = QPainter(image)
        doc.drawContents(painter, QRectF(image.rect())) 
        painter.end()
        image.save(fileName)
        if show_msg:
            QMessageBox.information(self, "Success", f"Image saved to:\n{fileName}")
        return fileName

    def _perform_word_ppt_save(self, fileName=None, file_type='word', show_msg=False):
        
        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        party_name = self.parent().party_name.text() if hasattr(self.parent(), 'party_name') else "N/A"
        
        if file_type == 'word':
            filter_str, default_ext, title_verb = "Word Documents (*.docx)", ".docx", "Document"
        else: # ppt
            filter_str, default_ext, title_verb = "PowerPoint Presentations (*.pptx)", ".pptx", "Presentation"
        
        if fileName is None:
            title = f"Save {file_type.title()} for Sharing"
            fileName, _ = QFileDialog.getSaveFileName(
                self, title, f"Order_{order_no}{default_ext}", filter_str
            )
            if not fileName: return None
            show_msg = True

        try:
            if file_type == 'word':
                from docx import Document    
                doc = Document()
                doc.add_heading('Order Report', 0)
                doc.add_paragraph(f'Order No: {order_no}')
                doc.add_paragraph(f'Party Name: {party_name}')
                doc.add_heading('Item Details:', level=2)
                doc.add_paragraph(self.content_data) 
                
                doc.save(fileName)

            elif file_type == 'ppt':
                from pptx import Presentation
                from pptx.util import Inches
                prs = Presentation()
                slide_layout = prs.slide_layouts[0]
                slide = prs.slides.add_slide(slide_layout)
                
                title = slide.shapes.title
                title.text = "Order Report"
                
                body = slide.placeholders[1]
                body.text = f"Order No: {order_no}\nParty Name: {party_name}\n\nItem Details:\n{self.content_data}"               
                prs.save(fileName)
            if show_msg:
                QMessageBox.information(self, "Success", f"{title_verb} saved to:\n{fileName}")
            return fileName
        except ImportError:
            QMessageBox.critical(self, "Error", f"The required library for {file_type.title()} export is missing. Please install 'python-docx' or 'python-pptx'.")
            return None
        except Exception as e:
            QMessageBox.critical(self, "Error", f"An error occurred during {file_type.title()} export: {e}")
            return None
        
    def share_via_whatsapp(self, file_path=None):
        """
        Opens WhatsApp Web/Desktop with a pre-filled message.
        Requires manual attachment of the file saved at file_path.
        """
        if file_path is None:
            QMessageBox.warning(self, "Error", "Please export the file first.")
            return

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        
        # Determine file type for the message
        ext = file_path.split('.')[-1].upper()
        
        message = f"Please find the Order Report (Order No: {order_no}) in {ext} format."
        
        encoded_message = QUrl.toPercentEncoding(message).data().decode()
        # NOTE: Using web.whatsapp.com for web/desktop sharing
        url = f"https://web.whatsapp.com/send?text={encoded_message}" 
        
        # Open the browser
        webbrowser.open(url)
        QMessageBox.information(self, "Manual Step Required", 
                                f"Your browser has opened WhatsApp. Please **manually attach** the saved file:\n\n{file_path}")

    def share_office_format_via_whatsapp(self, file_type):
        """Saves and shares Excel or Word/PPT via WhatsApp."""
        file_path = None
        
        if file_type == 'excel':
            file_path = self._perform_excel_save(None) # Pass None to prompt file dialog
        elif file_type == 'word':
            file_path = self._perform_word_ppt_save(None, 'word')
        
        if file_path:
            self.share_via_whatsapp(file_path)
            
    def share_via_whatsapp(self, file_path=None):
        """Opens WhatsApp link, prompting user to manually attach the file."""
        if file_path is None:
            QMessageBox.warning(self, "Error", "Please export the file first.")
            return

        order_no = self.parent().order_number.text() if hasattr(self.parent(), 'order_number') else "N/A"
        ext = file_path.split('.')[-1].upper()
        message = f"Please find the Order Report (Order No: {order_no}) in {ext} format."
        encoded_message = QUrl.toPercentEncoding(message).data().decode()
        url = f"https://web.whatsapp.com/send?text={encoded_message}" 
        # 1. Open the WhatsApp link (opens browser/desktop app)
        webbrowser.open(url)

        # 2. Open the file's containing folder (Cross-platform way to show the file)
        try:
            if sys.platform == 'win32':
                # Windows: Opens Explorer and selects the file
                os.startfile(os.path.dirname(file_path)) 
            elif sys.platform == 'darwin':
                # macOS: Opens Finder and selects the file
                os.system(f'open -R "{file_path}"')
            else:
                # Linux: Opens the file's directory
                os.system(f'xdg-open "{os.path.dirname(file_path)}"')
                
            QMessageBox.information(
                self, 
                "Action Required 🚨", 
                f"1. WhatsApp has opened in your browser/desktop app.\n"
                f"2. The file folder has also opened.\n\n"
                f"Please **drag and drop** the file onto the chat window to share it!"
            )
        except Exception as e:
            QMessageBox.warning(
                self, 
                "Action Required 🚨", 
                f"WhatsApp link opened. Could not open file explorer automatically.\n\n"
                f"Please manually navigate to and attach the saved file:\n\n{file_path}"
            )

    def show_whatsapp_share_menu_from_preview(self):
        """Shows format options for sharing, using the reusable save functions."""
        menu = QMenu(self)
        
        pdf_action = menu.addAction("Share as PDF")
        pdf_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_pdf_save))
        
        image_action = menu.addAction("Share as Image (PNG/JPG)")
        image_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_image_save))
        
        menu.addSeparator()
        excel_action = menu.addAction("Share as Excel (.xlsx)")
        excel_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_excel_save))
        
        word_action = menu.addAction("Share as Word (.docx)")
        word_action.triggered.connect(lambda: self.share_file_and_whatsapp(self._perform_word_ppt_save, 'word'))

        menu.exec_(QCursor.pos())

    def share_file_and_whatsapp(self, save_func, file_type=None):
        
        if file_type:
            file_path = save_func(None, file_type)
        else:
            file_path = save_func(None)
            
        if file_path:
            self.share_via_whatsapp(file_path)

class QuotationPreviewDialog(QDialog, ExportShareMixin):
    def __init__(self, parent, content_data, **kwargs):
        QDialog.__init__(self, parent, **kwargs)
        self.setWindowTitle("Quotation")
        self.content_data = content_data 
        self.document_type = "QUOTATION"
        self.setGeometry(200, 200, 350, 150) # Smaller dialog for simple actions

        main_layout = QVBoxLayout(self)

        print_layout = QHBoxLayout()
        direct_print_btn = QPushButton("🖨 Direct Print")
        preview_btn = QPushButton("🔍 Print Preview")
        
        print_layout.addWidget(direct_print_btn)
        print_layout.addWidget(preview_btn)
        main_layout.addLayout(print_layout)

        direct_print_btn.clicked.connect(self.direct_print)
        preview_btn.clicked.connect(self.show_preview)
        
    def _clean_table_html(self, html_content):
        if not html_content:
            return ""
        def clean_header(match):
            header_row = match.group(0)
            pattern = r'(Total Price<\/th>)'
            
            match_end = re.search(pattern, header_row, re.IGNORECASE)
            if match_end:
                return header_row[:match_end.end()] + '</tr>'
            return header_row # Return original if Total Price not found
        def clean_data_row(row_match):
            data_row = row_match.group(0)
            
            td_matches = list(re.finditer(r'<td.*?>.*?<\/td>', data_row, re.DOTALL | re.IGNORECASE))
            
            if len(td_matches) >= 7:
                total_price_td_end = td_matches[6].end()
                tr_end_match = re.search(r'<\/tr>', data_row)
                if tr_end_match:
                    tr_end_pos = tr_end_match.start()
                    if total_price_td_end < tr_end_pos:
                        return data_row[:total_price_td_end] + data_row[tr_end_pos:]
            
            return data_row # Return original if criteria not met

        rows = re.split(r'(<tr>.*?<\/tr>)', html_content, flags=re.DOTALL | re.IGNORECASE)
        
        cleaned_rows = []
        if rows:
            is_header_found = False
            for i, row_part in enumerate(rows):
                if not row_part.strip():
                    continue

                if '<tr>' in row_part and '</tr>' in row_part:
                    # This is a full row block
                    if not is_header_found:
                        cleaned_rows.append(clean_header(re.search(r'<tr>.*?<\/tr>', row_part, re.DOTALL | re.IGNORECASE)))
                        is_header_found = True
                    else:
                        cleaned_rows.append(clean_data_row(re.search(r'<tr>.*?<\/tr>', row_part, re.DOTALL | re.IGNORECASE)))
                else: # Any text before/after the table (like wrapper divs)
                    cleaned_rows.append(row_part)
        return "".join(cleaned_rows).replace('None', '') # Final cleanup and reassembly
    
    def _get_parent_text(self, attribute_name, default="N/A"):
        parent = self.parent()
        if hasattr(parent, attribute_name) and getattr(parent, attribute_name):
            widget = getattr(parent, attribute_name)
            if hasattr(widget, 'text'):
                return widget.text()
            elif hasattr(widget, 'currentText'):
                return widget.currentText() # For combo boxes
        return default

    def _get_parent_checkbox_state(self, checkbox_attr_name, price_attr_name=None):
        parent = self.parent()
        if hasattr(parent, checkbox_attr_name):
            widget = getattr(parent, checkbox_attr_name)
            if hasattr(widget, 'isChecked') and widget.isChecked():
                option_text = widget.text().strip()
                if price_attr_name:
                    price = self._get_parent_text(price_attr_name, "0")
                    return f"{option_text} (Price: {price})"
                return option_text
        return None

    def _get_parent_printing_options(self):
        parent = self.parent()
        printing_options = []
        keys = ['front', 'back', 'patch', 'embroidery', 'Dtf', 'Front sablimation', 'Back sablimation']
        
        if hasattr(parent, 'print_vars'):
            print_vars = getattr(parent, 'print_vars')
            for key in keys:
                if key in print_vars:
                    checkbox, price_edit = print_vars[key]
                    if checkbox.isChecked():
                        price = price_edit.text() if hasattr(price_edit, 'text') else "0"
                        option_text = checkbox.text().strip()
                        printing_options.append(f"<li>{option_text} (Price: {price} INR)</li>")
        return "".join(printing_options) or "<li>None Selected</li>"
    
    def _get_parent_track_options(self):
        parent = self.parent()
        track_options_list = []
        
        if hasattr(parent, 'track_vars') and hasattr(parent, 'track_extra_vars'):
            track_vars = getattr(parent, 'track_vars')
            track_extra_vars = getattr(parent, 'track_extra_vars')
            keys = ['Dori', '1 Piping', '2 Piping', 'Other'] 
            
            for key in keys:
                if key in track_vars:
                    checkbox, price_edit = track_vars[key]
                    extra_edit = track_extra_vars.get(key)
                    
                    if checkbox.isChecked():
                        price = price_edit.text() if hasattr(price_edit, 'text') else "0"
                        extra_detail = extra_edit.text() if extra_edit and hasattr(extra_edit, 'text') else ""
                        option_text = checkbox.text().strip()
                        detail_info = f" ({extra_detail})" if extra_detail else ""
                        track_options_list.append(f"<li>{option_text}{detail_info} (Price: {price} INR)</li>")
        return "".join(track_options_list) or "<li>None Selected</li>"

    def _get_tax_info(self, main_window):
        try:
            total_items_price = main_window._total_items_sum
            tax_percent = main_window._tax_percentage
            tax_amount = main_window._tax_amount
            grand_total = main_window._grand_total # This is the final, calculated number
        except AttributeError:
            # Fallback if the attributes haven't been calculated yet (shouldn't happen if called correctly)
            total_items_price = 0.0
            tax_percent = 0.0
            tax_amount = 0.0
            grand_total = 0.0

        display_grand_total = f"{grand_total:.2f}"

        tax_summary_html = f"""
        <table style="width: 100%; text-align: right; border-collapse: collapse; margin-top: 10px;">
            <tr>
                <td style="padding: 3px 0; border-top: 1px solid #ddd;"><b>Total Items Price:</b></td>
                <td style="padding: 3px 0; border-top: 1px solid #ddd;">₹ {total_items_price:.2f}</td>
            </tr>
            <tr>
                <td style="padding: 3px 0;"><b>Tax (GST) @ {tax_percent:.1f}%:</b></td>
                <td style="padding: 3px 0;">₹ {tax_amount:.2f}</td>
            </tr>
            <tr>
                <td style="font-size: 14pt; color: #d9534f; padding: 5px 0; border-top: 2px solid #333; border-bottom: 2px solid #333;"><b>GRAND TOTAL:</b></td>
                <td style="font-size: 14pt; color: #d9534f; padding: 5px 0; border-top: 2px solid #333; border-bottom: 2px solid #333;">₹ {display_grand_total}</td>
            </tr>
        </table>
        """
        return tax_summary_html, display_grand_total

    # --- Print Methods ---
    def print_document(self, printer):
        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        doc.print_(printer)

    def direct_print(self):
        printer = QPrinter(QPrinter.HighResolution)
        dialog = QPrintDialog(printer, self)
        if dialog.exec_() == QDialog.Accepted:
            self.print_document(printer)

    def show_preview(self):
        printer = QPrinter(QPrinter.HighResolution)
        preview = QPrintPreviewDialog(printer, self)
        preview.paintRequested.connect(self.print_document)
        preview.exec_()

    def get_print_content(self):
        order_no = self._get_parent_text('order_number')
        party_name = self._get_parent_text('party_name')
        order_date = self._get_parent_text('order_date') 
        delivery_date = self._get_parent_text('delivery_date')
        address = self._get_parent_text('address') 
        remarks = self._get_parent_text('remark_input') 
        school_name = self._get_parent_text('school_name') 
        barcode = self._get_parent_text('barcode')

        parent = self.parent()
        tax_summary_html, grand_total = self._get_tax_info(parent)
        
        canvas_image_base64_uri = parent._capture_canvas_as_base64() if hasattr(parent, '_capture_canvas_as_base64') else ""
        reference_image_base64_uris = parent._get_reference_images_base64()
        collar_options_list = "".join(f"<li>{opt}</li>" for opt in [self._get_parent_checkbox_state('rb_self', 'collar_price_self'), self._get_parent_checkbox_state('rb_rib', 'collar_price_rib'), self._get_parent_checkbox_state('rb_patti', 'collar_price_patti')] if opt) or "<li>None Selected</li>"
        button_option_map = {'BUTTON': 'rb_button', 'PLAIN': 'rb_plain', 'BOX': 'rb_box', 'V+': 'rb_vplus'}
        button_options_list = "".join(f"<li>{self._get_parent_checkbox_state(attr)}</li>" for attr in button_option_map.values() if self._get_parent_checkbox_state(attr)) or "<li>None Selected (Default)</li>"
        printing_options_list = self._get_parent_printing_options()
        track_pant_options_list = self._get_parent_track_options()

        cleaned_content_data = self._clean_table_html(self.content_data)

        # --- NEW: Reference Image Block HTML ---
        reference_images_html = ""
        if reference_image_base64_uris:
            image_table_rows = ""
            num_images = len(reference_image_base64_uris)
            
            # Iterate through the list, grabbing two URIs at a time
            for i in range(0, num_images, 2):
                uri1 = reference_image_base64_uris[i]
                uri2 = reference_image_base64_uris[i+1] if i + 1 < num_images else None
                
                # Cell 1 (always present)
                cell1 = f"""
                    <td style="width: 50%; padding: 5 25px; text-align: center; vertical-align: top; border: none;">
                        <img src="{uri1}" 
                            width="230"
                            height="250"
                            style="max-width: 100%; max-height: 250px; 
                                    border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;"
                            alt="Reference Image 1"/>
                    </td>
                """
                
                # Cell 2 (optional)
                cell2 = ""
                if uri2:
                    cell2 = f"""
                        <td style="width: 50%; padding: 5 85px; text-align: center; vertical-align: top; border: none;">
                            <img src="{uri2}" 
                                width="230"
                                height="250"
                                style="max-width: 100%; max-height: 250px; 
                                        border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;"
                                alt="Reference Image 2"/>
                        </td>
                    """
                else:
                    cell2 = '<td style="width: 50%; border: none;"></td>' 
                
                image_table_rows += f'<tr style="page-break-inside: avoid;">{cell1}{cell2}</tr>'

            # Final HTML structure for the image block
            reference_images_html = f"""
            <div style="margin-top: 20px; clear: both;">
                <h2 class="section-header" style="background-color: #f0f0f0; border-left: 5px solid #007bff;">Customer Reference Images</h2>
                <table style="width: 100%; border-collapse: collapse; margin-top: 5px;">
                    <tbody>
                        {image_table_rows}
                    </tbody>
                </table>
            </div>
            """
        # --- END: FIX FOR REFERENCE IMAGES ---

        html_content = f"""
        <html>
        <head>
            <style>
                @page {{ size: A4; margin: 20mm; }} 
                body {{ font-family: 'Arial', sans-serif; font-size: 10pt; line-height: 1.4; }}
                h1 {{ text-align: center; margin-bottom: 5px; color: #333; }}
                hr {{ border: 0.5px solid #ccc; }}
                .company-header {{ text-align: center; margin-bottom: 20px; }}
                .company-header h2 {{ margin: 0; font-size: 14pt; color: #007bff; }}
                .company-header p {{ margin: 2px 0; font-size: 9pt; color: #555; }}
                .header-table {{ width: 100%; border-collapse: collapse; margin-bottom: 10px; }}
                .header-table td {{ padding: 3px 5px; vertical-align: top; }}
                .section-header {{ background-color: #f0f0f0; padding: 5px; margin-top: 15px; margin-bottom: 5px; border-left: 5px solid #007bff; font-size: 12pt; }}
                .main-content-footer {{display: flex; flex-direction: row; justify-content: space-between; align-items: flex-start; width: 100%;}}
                .item-details-section {{flex-grow: 1; flex-basis: 60%; min-width: 0%; padding-right: 15px; box-sizing: border-box; overflow: hidden;}}
                .item-table {{ width: 100%; border-collapse: collapse; margin-top: 10px; }}
                .item-table th, .item-table td {{ border: 1px solid #ddd; padding: 6px; text-align: left; }}
                .item-table th {{ background-color: #e9ecef; }}
                .summary-container {{flex-basis: 35%; width: 35%; padding: 0; text-align: right; margin-top: 15px;}}
                .product-design-preview {{ max-width: 150px; max-height: 250px; width: auto; height: auto; border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;}}
                .options-list {{list-style-type: disc; padding-left: 20px; margin: 0 0 10px 0; font-size: 10pt;}}
            </style>
        </head>
        <body>
            <div class="company-header">
                <h2>[YOUR COMPANY NAME HERE]</h2>
                <p>[YOUR ADDRESS LINE 1]</p>
                <p>[YOUR CONTACT INFO]</p>
                <hr>
            </div>
            <hr>

            <table class="header-table">
                <tr>
                    <td width="33%"><b>Order No:</b> {order_no}</td>
                    <td width="33%"><b>Order Date:</b> {order_date}</td>
                    <td width="34%"><b>Delivery Date:</b> {delivery_date}</td>
                </tr>
                <tr>
                    <td width="33%"><b>Barcode:</b> {barcode}</td>
                    <td colspan="2"><b>Address:</b> {address}</td>
                </tr>
                <tr>
                    <td width="33%"><b>Party Name:</b> {party_name}</td>
                    <td colspan="2"><b>School Name:</b> {school_name}</td>
                </tr>
            </table>

            <h2 class="section-header" style="margin-top: 10px; margin-bottom: 5px;">Product Design & Customization</h2>
            <table style="width: 100%; border-collapse: collapse; margin-top: 0; margin-bottom: 0;">
                <tr>
                    <td style="width: 50%; vertical-align: top; padding: 0 5px 0 0; text-align: center;">
                        <h3 style="margin-top: 0; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Product Image</h3>
                        <div style="max-width: 100%; margin: 0 auto; line-height: 1;">
                            {f'<img src="{canvas_image_base64_uri}" class="product-image-preview" alt="Product Design" style="max-height: 150px;"/>' if canvas_image_base64_uri else '<p style="margin: 0; font-size: 9pt;">No Product Image Selected</p>'}
                        </div>
                    </td>
                </tr>
                <tr>
                    <td style="width: 50%; vertical-align: top; border-left: 1px solid #ddd; padding: 0 0 0 10px;">
                        <h3 style="margin-top: 0; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Printing Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 3px;">{printing_options_list}</ul>
                        
                        <h3 style="margin-top: 3px; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Collar Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 3px;">{collar_options_list}</ul>
                        
                        <h3 style="margin-top: 3px; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Button Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 3px;">{button_options_list}</ul>

                        <h3 style="margin-top: 3px; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Track Pant Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 0;">{track_pant_options_list}</ul>
                    </td>
                </tr>
            </table>

            <h2 class="section-header">Item Details & Pricing</h2>
            
            <div class="main-content-footer">
                <div class="item-details-section">
                    <div class="item-table-container">
                        {cleaned_content_data}
                    </div>
                </div>

                <div class="summary-container">
                    {tax_summary_html}
                </div>
            </div>
            
            <div style="clear: both; margin-top: 10px;">
                <h2 class="section-header">Remark</h2>
                <p>{remarks if remarks != "N/A" and remarks else "No special remarks."}</p>
            </div>

            {reference_images_html}
            
            <div style="margin-top: 50px; text-align: center; font-size: 8pt; color: #777;">
                <p>Signature (Seller)</p>
            </div>

        </body>
        </html>
        """
        return html_content
    
    def print_document(self, printer):
        doc = QTextDocument()
        doc.setHtml(self.get_print_content())
        doc.print_(printer)

    def direct_print(self):
        printer = QPrinter(QPrinter.HighResolution)
        dialog = QPrintDialog(printer, self)
        if dialog.exec_() == QDialog.Accepted:
            self.print_document(printer)
    def show_preview(self):
        printer = QPrinter(QPrinter.HighResolution)
        preview = QPrintPreviewDialog(printer, self)
        preview.paintRequested.connect(self.print_document)

        export_btn = QPushButton("🔽 Save Options")
        export_btn.setToolTip("Save to PDF/Excel/Image/Word/PPT")
        export_btn.clicked.connect(self.show_export_menu)
        preview.layout().addWidget(export_btn) 

        share_btn = QPushButton("📱 Share via WhatsApp")
        share_btn.setToolTip("Share Order as PDF via WhatsApp")
        share_btn.clicked.connect(self.show_whatsapp_share_menu)
        preview.layout().addWidget(share_btn)

        preview.exec_()

class JobWorkPreviewDialog(QuotationPreviewDialog):
    
    def __init__(self, parent, content_data, **kwargs):
        super().__init__(parent, content_data, **kwargs)
        self.setWindowTitle("Job Work (Stretching)")
        self.document_type = "JOB_WORK"

    def _clean_table_html(self, html_content):
        
        if not html_content:
            return ""
        
        def clean_header(match):
            header_row = match.group(0)
            header_row = re.sub(r'<th[^>]*>Total Price<\/th>', '', header_row, flags=re.IGNORECASE)
            header_row = re.sub(r'<th[^>]*>Status<\/th>', '', header_row, flags=re.IGNORECASE)
            header_row = re.sub(r'<th[^>]*>Unit<\/th>', '', header_row, flags=re.IGNORECASE)
            
            return header_row
            
        def clean_data_row(row_match):
            data_row = row_match.group(0)
            td_matches = list(re.finditer(r'<td.*?<\/td>', data_row, re.DOTALL | re.IGNORECASE))

            if len(td_matches) >= 8:
                new_row_content = "".join(td_matches[i].group(0) for i in range(5))
                tr_end_match = re.search(r'<\/tr>', data_row)
                if tr_end_match:
                    return f'<tr>{new_row_content}{data_row[tr_end_match.start():]}'
            
            return data_row 

        rows = re.split(r'(<tr>.*?<\/tr>)', html_content, flags=re.DOTALL | re.IGNORECASE)
        
        cleaned_rows = []
        if rows:
            is_header_found = False
            for i, row_part in enumerate(rows):
                if not row_part.strip():
                    continue

                if '<tr>' in row_part and '</tr>' in row_part:
                    row_match = re.search(r'<tr>.*?<\/tr>', row_part, re.DOTALL | re.IGNORECASE)
                    if row_match:
                        if not is_header_found:
                            cleaned_rows.append(clean_header(row_match))
                            is_header_found = True
                        else:
                            cleaned_rows.append(clean_data_row(row_match))
                else: 
                    cleaned_rows.append(row_part)
        return "".join(cleaned_rows).replace('None', '')

    def _get_tax_info(self, main_window):
        return "", "0.00"

    def get_print_content(self):
        order_no = self._get_parent_text('order_number')
        barcode = self._get_parent_text('barcode')
        employee_name = self._get_parent_text('employee_name', 'N/A') # Assuming employee_name field exists on parent
        remarks = self._get_parent_text('remark_input') 

        parent = self.parent()
        tax_summary_html, grand_total = self._get_tax_info(parent) # Will be empty
        
        canvas_image_base64_uri = parent._capture_canvas_as_base64() if hasattr(parent, '_capture_canvas_as_base64') else ""
        reference_image_base64_uris = parent._get_reference_images_base64()
        
        collar_options_list = "".join(f"<li>{opt}</li>" for opt in [self._get_parent_checkbox_state('rb_self', 'collar_price_self'), self._get_parent_checkbox_state('rb_rib', 'collar_price_rib'), self._get_parent_checkbox_state('rb_patti', 'collar_price_patti')] if opt) or "<li>None Selected</li>"
        collar_options_list = re.sub(r' \(Price: .*?\)', '', collar_options_list) 
        
        button_option_map = {'BUTTON': 'rb_button', 'PLAIN': 'rb_plain', 'BOX': 'rb_box', 'V+': 'rb_vplus'}
        button_options_list = "".join(f"<li>{self._get_parent_checkbox_state(attr)}</li>" for attr in button_option_map.values() if self._get_parent_checkbox_state(attr)) or "<li>None Selected (Default)</li>"
        
        printing_options_list = self._get_parent_printing_options()
        printing_options_list = re.sub(r' \(Price: .*? INR\)', '', printing_options_list)

        track_pant_options_list = self._get_parent_track_options()
        track_pant_options_list = re.sub(r' \(Price: .*? INR\)', '', track_pant_options_list)

        cleaned_content_data = self._clean_table_html(self.content_data)

        reference_images_html = ""
        if reference_image_base64_uris:
            image_table_rows = ""
            num_images = len(reference_image_base64_uris)
            
            # Iterate through the list, grabbing two URIs at a time
            for i in range(0, num_images, 2):
                uri1 = reference_image_base64_uris[i]
                uri2 = reference_image_base64_uris[i+1] if i + 1 < num_images else None
                
                # Cell 1 (always present)
                cell1 = f"""
                    <td style="width: 50%; padding: 5 25px; text-align: center; vertical-align: top; border: none;">
                        <img src="{uri1}" 
                            width="230"
                            height="250"
                            style="max-width: 100%; max-height: 250px; 
                                    border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;"
                            alt="Reference Image 1"/>
                    </td>
                """
                
                # Cell 2 (optional)
                cell2 = ""
                if uri2:
                    cell2 = f"""
                        <td style="width: 50%; padding: 5 85px; text-align: center; vertical-align: top; border: none;">
                            <img src="{uri2}" 
                                width="230"
                                height="250"
                                style="max-width: 100%; max-height: 250px; 
                                        border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;"
                                alt="Reference Image 2"/>
                        </td>
                    """
                else:
                    cell2 = '<td style="width: 50%; border: none;"></td>' 
                
                image_table_rows += f'<tr style="page-break-inside: avoid;">{cell1}{cell2}</tr>'

            # Final HTML structure for the image block
            reference_images_html = f"""
            <div style="margin-top: 20px; clear: both;">
                <h2 class="section-header" style="background-color: #f0f0f0; border-left: 5px solid #007bff;">Customer Reference Images</h2>
                <table style="width: 100%; border-collapse: collapse; margin-top: 5px;">
                    <tbody>
                        {image_table_rows}
                    </tbody>
                </table>
            </div>
            """
        # --- END: FIX FOR REFERENCE IMAGES ---

        html_content = f"""
        <html>
        <head>
            <style>
                @page {{ size: A4; margin: 20mm; }} 
                body {{ font-family: 'Arial', sans-serif; font-size: 10pt; line-height: 1.4; }}
                h1 {{ text-align: center; margin-bottom: 5px; color: #333; }}
                hr {{ border: 0.5px solid #ccc; }}
                .company-header {{ text-align: center; margin-bottom: 20px; }}
                .company-header h2 {{ margin: 0; font-size: 16pt; color: #d9534f; }}
                .company-header p {{ margin: 2px 0; font-size: 9pt; color: #555; }}
                .header-table {{ width: 100%; border-collapse: collapse; margin-bottom: 10px; }}
                .header-table td {{ padding: 3px 5px; vertical-align: top; }}
                .section-header {{ background-color: #f0f0f0; padding: 5px; margin-top: 15px; margin-bottom: 5px; border-left: 5px solid #007bff; font-size: 12pt; }}
                .item-table-container {{ overflow-x: auto; }}
                .item-table {{ width: 100%; border-collapse: collapse; margin-top: 10px; }}
                .item-table th, .item-table td {{ border: 1px solid #ddd; padding: 6px; text-align: left; }}
                .item-table th {{ background-color: #e9ecef; }}
                .product-design-preview {{ max-width: 150px; max-height: 250px; width: auto; height: auto; border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;}}
                .options-list {{list-style-type: disc; padding-left: 20px; margin: 0 0 10px 0; font-size: 10pt;}}
            </style>
        </head>
        <body>
            <div class="company-header">
                <h2>JOB WORK (STRETCHING) SLIP</h2>
                <p>Company Logo and Company Details</p>
                <hr>
            </div>

            <table class="header-table">
                <tr>
                    <td width="33%"><b>Order No:</b> {order_no}</td>
                    <td width="33%"><b>Barcode:</b> {barcode}</td>
                    <td width="34%"><b>Employee Name:</b> {employee_name}</td>
                </tr>
            </table>

            <h2 class="section-header" style="margin-top: 10px; margin-bottom: 5px;">Product Design & Customization</h2>
            <table style="width: 100%; border-collapse: collapse; margin-top: 0; margin-bottom: 0;">
                <tr>
                    <td style="width: 50%; vertical-align: top; padding: 0 5px 0 0; text-align: center;">
                        <h3 style="margin-top: 0; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Product Image</h3>
                        <div style="max-width: 100%; margin: 0 auto; line-height: 1;">
                            {f'<img src="{canvas_image_base64_uri}" class="product-design-preview" alt="Product Design" style="max-height: 150px;"/>' if canvas_image_base64_uri else '<p style="margin: 0; font-size: 9pt;">No Product Image</p>'}
                        </div>
                    </td>
                </tr>
                <tr>    
                    <td style="width: 50%; vertical-align: top; border-left: 1px solid #ddd; padding: 0 0 0 10px;">
                        <h3 style="margin-top: 0; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Printing Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 3px;">{printing_options_list}</ul>
                        
                        <h3 style="margin-top: 3px; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Collar Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 3px;">{collar_options_list}</ul>
                        
                        <h3 style="margin-top: 3px; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Button Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 3px;">{button_options_list}</ul>

                        <h3 style="margin-top: 3px; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Track Pant Options (Only Selected)</h3>
                        <ul class="options-list" style="margin-bottom: 0;">{track_pant_options_list}</ul>
                    </td>
                </tr>
            </table>

            <h2 class="section-header">Item Details</h2>
            <div class="item-table-container">
                {cleaned_content_data}
            </div>

            {tax_summary_html} <div style="clear: both; margin-top: 10px;">
                <h2 class="section-header">Remark</h2>
                <p>{remarks if remarks != "N/A" and remarks else "No special remarks."}</p>
            </div>

            {reference_images_html}
            
            <div style="margin-top: 50px; text-align: center; font-size: 8pt; color: #777;">
                <p>Signature (Job Work Manager)</p>
            </div>

        </body>
        </html>
        """
        return html_content
    
class CuttingJobPreviewDialog(JobWorkPreviewDialog):

    def __init__(self, parent=None, html_content=""):
        super().__init__(parent, html_content)
        self.setWindowTitle("Cutting Job Slip Preview")

    def get_print_content(self):
        current_date = QDate.currentDate().toString("dd-MM-yyyy")

        base_html = super().get_print_content()

        order_no = self._get_parent_text('order_number')
        barcode = self._get_parent_text('barcode')
        employee_name = self._get_parent_text('employee_name', 'N/A')
        
        new_header_table_html = f"""
            <table class="header-table">
                <tr>
                    <td width="33%"><b>Order No:</b> {order_no}</td>
                    <td width="33%"><b>Barcode:</b> {barcode}</td>
                    <td width="34%"><b>Current Date:</b> {current_date}</td>
                </tr>
                <tr>
                    <td width="33%"><b>Employee Name:</b> {employee_name}</td>
                    <td colspan="2"></td> 
                </tr>
            </table>
        """
        new_html = re.sub(
            r'<table class="header-table">.*?<\/table>', 
            new_header_table_html, 
            base_html, 
            flags=re.DOTALL | re.IGNORECASE,
            count=1 
        )
        new_html = new_html.replace("JOB WORK (STRETCHING) SLIP", "CUTTING JOB SLIP")
        
        images_block_regex = r'<div\s+style="margin-top:\s*20px;.*?">.*?Customer Reference Images<\/h2>.*?<\/div>'
        final_html = re.sub(
            images_block_regex, 
            '', # Replace with empty string
            new_html, 
            flags=re.DOTALL | re.IGNORECASE
        )
        return final_html

class PrintingJobPreviewDialog(JobWorkPreviewDialog):

    def __init__(self, parent, content_data, **kwargs):
        super().__init__(parent, content_data, **kwargs)
        self.setWindowTitle("Printing Job Slip")
        self.document_type = "PRINTING_JOB"
        
    def _clean_table_html(self, html_content):

        if not html_content:
            return ""
        
    def get_print_content(self):
        order_no = self._get_parent_text('order_number')
        barcode = self._get_parent_text('barcode')
        current_date = self._get_parent_text('order_date')
        employee_name = self._get_parent_text('employee_name', 'N/A') 
        school_name = self._get_parent_text('school_name', 'N/A')
        remarks = self._get_parent_text('remark_input') 

        parent = self.parent()
        canvas_image_base64_uri = parent._capture_canvas_as_base64() if hasattr(parent, '_capture_canvas_as_base64') else ""
        html_content = f"""
        <html>
        <head>
            <style>
                @page {{ size: A4; margin: 20mm; }} 
                body {{ font-family: 'Arial', sans-serif; font-size: 10pt; line-height: 1.4; }}
                h1 {{ text-align: center; margin-bottom: 5px; color: #333; }}
                hr {{ border: 0.5px solid #ccc; }}
                .company-header {{ text-align: center; margin-bottom: 20px; }}
                .company-header h2 {{ margin: 0; font-size: 16pt; color: #d9534f; }}
                .company-header p {{ margin: 2px 0; font-size: 9pt; color: #555; }}
                .header-table {{ width: 100%; border-collapse: collapse; margin-bottom: 10px; }}
                .header-table td {{ padding: 3px 5px; vertical-align: top; }}
                .section-header {{ background-color: #f0f0f0; padding: 5px; margin-top: 15px; margin-bottom: 5px; border-left: 5px solid #007bff; font-size: 12pt; }}
                .item-table-container {{ overflow-x: auto; }}
                .item-table {{ width: 100%; border-collapse: collapse; margin-top: 10px; }}
                .item-table th, .item-table td {{ border: 1px solid #ddd; padding: 6px; text-align: left; }}
                .item-table th {{ background-color: #e9ecef; }}
                .product-design-preview {{ max-width: 150px; max-height: 250px; width: auto; height: auto; border: 1px solid #ccc; object-fit: contain; display: block; margin: 0 auto;}}
                .options-list {{list-style-type: disc; padding-left: 20px; margin: 0 0 10px 0; font-size: 10pt;}}
            </style>
        </head>
        <body>
            <div class="company-header">
                <h2>PRINTING JOB SLIP</h2>
                <p>Company Logo and Company Details</p>
                <hr>
            </div>

            <table class="header-table">
                <tr>
                    <td width="33%"><b>Order No:</b> {order_no}</td>
                    <td width="33%"><b>Barcode:</b> {barcode}</td>
                    <td width="34%"><b>Current Date:</b> {current_date}</td>
                </tr>
                <tr>
                    <td width="33%"><b>Employee Name:</b> {employee_name}</td>
                    <td colspan="2"><b>School Name:</b> {school_name}</td>
                </tr>
            </table>

            <h2 class="section-header" style="margin-top: 10px; margin-bottom: 5px;">Product Design & Customization Details</h2>
            <table style="width: 100%; border-collapse: collapse; margin-top: 0; margin-bottom: 0;">
                <tr>
                    <td style="width: 50%; vertical-align: top; padding: 0 5px 0 0; text-align: center;">
                        <h3 style="margin-top: 0; margin-bottom: 3px; font-size: 11pt; color: #007bff;">Product Image</h3>
                        <div style="max-width: 100%; margin: 0 auto; line-height: 1;">
                            {f'<img src="{canvas_image_base64_uri}" class="product-design-preview" alt="Product Design" style="max-height: 150px;"/>' if canvas_image_base64_uri else '<p style="margin: 0; font-size: 9pt;">No Product Image</p>'}
                        </div>
                    </td>
                </tr>
            </table>

            <div style="clear: both; margin-top: 10px;">
                <h2 class="section-header">Remark</h2>
                <p>{remarks if remarks != "N/A" and remarks else "No special remarks."}</p>
            </div>
            
            <div style="margin-top: 50px; text-align: center; font-size: 8pt; color: #777;">
                <p>Signature (Printing Manager)</p>
            </div>

        </body>
        </html>
        """
        return html_content